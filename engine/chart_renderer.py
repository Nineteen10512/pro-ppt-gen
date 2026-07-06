"""Semantic chart renderer – native OOXML chart via python-pptx.

This module is the ONLY place that touches python-pptx chart APIs.
It consumes a semantic ``chart_data`` dict (produced by LLM/layouts) and
a bounding box (x, y, cx, cy) computed by the grid system, then renders
an editable chart with theme-aware colours, fonts and styling.

Semantic type -> XL_CHART_TYPE mapping (single source of truth):
    column            -> COLUMN_CLUSTERED
    bar               -> BAR_CLUSTERED
    stacked_column    -> COLUMN_STACKED
    stacked_bar       -> BAR_STACKED
    line              -> LINE
    line_markers      -> LINE_MARKERS
    pie               -> PIE
    doughnut          -> DOUGHNUT
    area              -> AREA
    scatter           -> XY_SCATTER
    radar             -> RADAR
"""
from __future__ import annotations

from typing import Any

from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LEGEND_POSITION,
    XL_TICK_MARK,
)
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Semantic type -> XL_CHART_TYPE
# ---------------------------------------------------------------------------
CHART_TYPE_MAP: dict[str, XL_CHART_TYPE] = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
    "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "radar": XL_CHART_TYPE.RADAR,
}

# Chart types that are category-based (as opposed to XY scatter).
_CATEGORY_TYPES = {
    "column", "bar", "stacked_column", "stacked_bar",
    "line", "line_markers", "pie", "doughnut", "area", "radar",
}
_SCATTER_TYPES = {"scatter"}
# Chart types that don't use axes / gridlines.
_NO_AXIS_TYPES = {"pie", "doughnut"}
# Chart types where data labels default ON.
_DEFAULT_LABELS_TYPES = {"pie", "doughnut"}

_LEGEND_POS_MAP = {
    "bottom": XL_LEGEND_POSITION.BOTTOM,
    "top": XL_LEGEND_POSITION.TOP,
    "left": XL_LEGEND_POSITION.LEFT,
    "right": XL_LEGEND_POSITION.RIGHT,
}


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def _set_east_asia_font(rPr, font_name: str):
    """Append an <a:ea typeface="..."/> element to a run/paragraph properties element."""
    for ea in rPr.findall(qn("a:ea")):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_name)


def _style_font(font, tokens: dict, size_key: str = "label_font_size",
                color: RGBColor | None = None, bold: bool = False,
                cn_font: str = "微软雅黑"):
    """Apply theme-aware font styling: latin + eastAsia + size + color."""
    size = tokens["chart"][size_key]
    font.size = size
    font.bold = bold
    font.name = tokens["font"]["family"]["body"]  # latin
    if color is not None:
        font.color.rgb = color
    # Set eastAsian font on the underlying defRPr / txPr
    rPr = None
    for attr in ("_rPr",):
        rPr = getattr(font, attr, None)
        if rPr is not None:
            break
    if rPr is None:
        # Fall back to XML walk
        el = font._element  # CT_TextFont
        rPr = el.find(qn("a:defRPr")) if el is not None else None
    if rPr is None:
        # Most pptx chart fonts: go through the font._element's parent chain
        try:
            rPr = font._r.get_or_add_rPr()
        except Exception:
            rPr = None
    if rPr is not None:
        _set_east_asia_font(rPr, cn_font)


def _set_text_frame_cjk(tf, cn_font: str):
    """Walk every paragraph/run in a chart text frame and set eastAsia font."""
    for p in tf.paragraphs:
        for r in p.runs:
            try:
                rPr = r._r.get_or_add_rPr()
                _set_east_asia_font(rPr, cn_font)
            except Exception:
                pass


def _set_axis_cjk_font(axis, cn_font: str, tokens: dict):
    """Set axis tick label font + eastAsia."""
    try:
        tf = axis.tick_labels.font
        axis.tick_labels.font.size = tokens["chart"]["label_font_size"]
        axis.tick_labels.font.name = tokens["font"]["family"]["body"]
        axis.tick_labels.font.color.rgb = tokens["color"]["text"]
    except Exception:
        pass
    # axis title if present
    try:
        if axis.has_title:
            _set_text_frame_cjk(axis.axis_title.text_frame, cn_font)
    except Exception:
        pass


def _set_legend_cjk(legend, cn_font: str, tokens: dict):
    try:
        legend.font.size = tokens["chart"]["legend_font_size"]
        legend.font.name = tokens["font"]["family"]["body"]
        legend.font.color.rgb = tokens["color"]["text"]
        # CJK via XML
        tf_elem = legend._element  # CT_Legend
        # Find all run properties inside legend entries and set eastAsia
        for rPr in tf_elem.iter(qn("a:defRPr")):
            _set_east_asia_font(rPr, cn_font)
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Build chart data
# ---------------------------------------------------------------------------
def _build_category_data(chart_data: dict) -> CategoryChartData:
    data = CategoryChartData()
    data.categories = list(chart_data["categories"])
    for s in chart_data["series"]:
        data.add_series(s["name"], tuple(s["values"]))
    return data


def _build_xy_data(chart_data: dict) -> XyChartData:
    data = XyChartData()
    for s in chart_data["series"]:
        series_data = data.add_series(s["name"])
        for point in s["values"]:
            if isinstance(point, dict):
                series_data.add_data_point(point.get("x", 0), point.get("y", 0))
            else:
                # Fallback: treat as (index, value)
                series_data.add_data_point(len(series_data.x_values), point)
    return data


# ---------------------------------------------------------------------------
# Post-style the chart element
# ---------------------------------------------------------------------------
def _style_chart(chart, chart_data: dict, tokens: dict, xl_type: XL_CHART_TYPE, cn_font: str):
    palette = tokens["color"]["chart_palette"]
    gridline_color = tokens["color"]["chart_gridline"]
    text_color = tokens["color"]["text"]
    axis_size = tokens["chart"]["label_font_size"]
    data_label_size = tokens["chart"]["data_label_font_size"]
    marker_size = tokens["chart"]["marker_size"]
    show_legend = bool(chart_data.get("show_legend", True))
    legend_pos = chart_data.get("legend_position", "bottom")
    show_data_labels = bool(chart_data.get(
        "show_data_labels", chart_data["type"] in _DEFAULT_LABELS_TYPES
    ))
    number_format = chart_data.get("number_format", "0")
    is_pie_like = chart_data["type"] in _NO_AXIS_TYPES
    y_axis_cfg = chart_data.get("y_axis") or {}

    # ---------- Transparent backgrounds ----------
    # chart area no fill / no line
    try:
        chart.has_title = bool(chart_data.get("title"))
    except Exception:
        pass
    try:
        el = chart._chartSpace
        # No rounded corners / no shadow
        for spPr in el.iter(qn("c:spPr")):
            # no fill
            noFill = etree.SubElement(spPr, qn("a:noFill")) if spPr.find(qn("a:noFill")) is None else spPr.find(qn("a:noFill"))
            # remove existing solidFill children
            for sf in spPr.findall(qn("a:solidFill")):
                spPr.remove(sf)
        # plot area transparent
        plotArea = el.find(qn("c:plotArea"))
        if plotArea is not None:
            spPr = plotArea.find(qn("c:spPr"))
            if spPr is None:
                spPr = etree.SubElement(plotArea, qn("c:spPr"))
                plotArea.insert(0, spPr)
            for sf in spPr.findall(qn("a:solidFill")):
                spPr.remove(sf)
            if spPr.find(qn("a:noFill")) is None:
                etree.SubElement(spPr, qn("a:noFill"))
            ln = spPr.find(qn("a:ln"))
            if ln is None:
                ln = etree.SubElement(spPr, qn("a:ln"))
            if ln.find(qn("a:noFill")) is None:
                for sf in ln.findall(qn("a:solidFill")):
                    ln.remove(sf)
                etree.SubElement(ln, qn("a:noFill"))
    except Exception:
        pass

    # ---------- Title ----------
    if chart.has_title and chart_data.get("title"):
        try:
            chart.chart_title.text_frame.text = chart_data["title"]
            tf = chart.chart_title.text_frame
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.size = tokens["chart"]["title_font_size"]
                    r.font.bold = True
                    r.font.name = tokens["font"]["family"]["heading"]
                    r.font.color.rgb = text_color
                    try:
                        rPr = r._r.get_or_add_rPr()
                        _set_east_asia_font(rPr, cn_font)
                    except Exception:
                        pass
        except Exception:
            pass

    # ---------- Legend ----------
    chart.has_legend = show_legend
    if show_legend:
        legend = chart.legend
        legend.include_in_layout = False
        pos_enum = _LEGEND_POS_MAP.get(legend_pos, XL_LEGEND_POSITION.BOTTOM)
        try:
            legend.position = pos_enum
        except Exception:
            pass
        _set_legend_cjk(legend, cn_font, tokens)

    # ---------- Series colors ----------
    try:
        for idx, series in enumerate(chart.series):
            color = palette[idx % len(palette)]
            # Determine the shape XML based on chart type
            ser_elem = series._element  # CT_BarSer etc.
            spPr = ser_elem.find(qn("c:spPr"))
            if spPr is None:
                # Insert spPr before tx (marker position depending)
                spPr = etree.SubElement(ser_elem, qn("c:spPr"))
                # Move to near top (after idx/order/tx)
                # Find insertion point: after <c:tx>
                tx = ser_elem.find(qn("c:tx"))
                if tx is not None:
                    ser_elem.remove(spPr)
                    tx.addnext(spPr)
            # Clean existing fill
            for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:pattFill"):
                for el in spPr.findall(qn(tag)):
                    spPr.remove(el)
            solid = etree.SubElement(spPr, qn("a:solidFill"))
            srgb = etree.SubElement(solid, qn("a:srgbClr"))
            srgb.set("val", "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2]))

            # Line thickness for line-type charts
            ctype = chart_data["type"]
            if ctype in ("line", "line_markers", "radar", "scatter", "area"):
                ln = spPr.find(qn("a:ln"))
                if ln is None:
                    ln = etree.SubElement(spPr, qn("a:ln"))
                ln.set("w", "22225")  # ~1.75pt line weight
                for sf in ln.findall(qn("a:solidFill")):
                    ln.remove(sf)
                ln_solid = etree.SubElement(ln, qn("a:solidFill"))
                ln_srgb = etree.SubElement(ln_solid, qn("a:srgbClr"))
                ln_srgb.set("val", "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2]))

            # Markers for line_markers / scatter
            if ctype in ("line_markers", "scatter"):
                marker = ser_elem.find(qn("c:marker"))
                if marker is None:
                    marker = etree.SubElement(ser_elem, qn("c:marker"))
                for sym in marker.findall(qn("c:symbol")):
                    marker.remove(sym)
                sym = etree.SubElement(marker, qn("c:symbol"))
                sym.set("val", "circle")
                for sz in marker.findall(qn("c:size")):
                    marker.remove(sz)
                sz_el = etree.SubElement(marker, qn("c:size"))
                sz_el.set("val", str(marker_size))
                mSpPr = marker.find(qn("c:spPr"))
                if mSpPr is None:
                    mSpPr = etree.SubElement(marker, qn("c:spPr"))
                for sf in mSpPr.findall(qn("a:solidFill")):
                    mSpPr.remove(sf)
                m_solid = etree.SubElement(mSpPr, qn("a:solidFill"))
                m_srgb = etree.SubElement(m_solid, qn("a:srgbClr"))
                m_srgb.set("val", "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2]))

            # Data labels
            if show_data_labels:
                _enable_data_labels(series, number_format, ctype, data_label_size,
                                    text_color, color, cn_font, tokens)
    except Exception as e:
        # Non-fatal
        pass

    # ---------- Axes & gridlines (non-pie) ----------
    if not is_pie_like:
        try:
            category_axis = None
            value_axis = None
            for ax in [getattr(chart, "category_axis", None), getattr(chart, "value_axis", None)]:
                pass
            # python-pptx exposes axis attrs per chart type; iterate via plot
            plot = chart.plots[0]
            # Try to get axes
            c_ax = getattr(plot, "category_axis", None) or getattr(chart, "category_axis", None)
            v_ax = getattr(plot, "value_axis", None) or getattr(chart, "value_axis", None)

            def _style_axis(ax, is_value: bool):
                if ax is None:
                    return
                try:
                    ax.tick_labels.font.size = axis_size
                    ax.tick_labels.font.name = tokens["font"]["family"]["body"]
                    ax.tick_labels.font.color.rgb = text_color
                    # Tick marks outward
                    try:
                        ax.major_tick_mark = XL_TICK_MARK.OUTSIDE
                    except Exception:
                        pass
                    try:
                        ax.minor_tick_mark = XL_TICK_MARK.NONE
                    except Exception:
                        pass
                    # Axis line colour = text_light
                    try:
                        ax.format.line.color.rgb = tokens["color"]["text_light"]
                    except Exception:
                        pass
                    # Major gridlines
                    if is_value:
                        try:
                            gl = ax.major_gridlines
                            gl.format.line.color.rgb = gridline_color
                            gl.format.line.width = Pt(tokens["chart"]["gridline_width_pt"])
                            # Dashed
                            ln_pr = gl.format.line._get_or_add_ln()
                            prstDash = ln_pr.find(qn("a:prstDash"))
                            if prstDash is None:
                                prstDash = etree.SubElement(ln_pr, qn("a:prstDash"))
                            prstDash.set("val", "dash")
                        except Exception:
                            pass
                        # Y axis min/max
                        if y_axis_cfg.get("min") is not None:
                            try:
                                ax.minimum_scale = y_axis_cfg["min"]
                            except Exception:
                                pass
                        if y_axis_cfg.get("max") is not None:
                            try:
                                ax.maximum_scale = y_axis_cfg["max"]
                            except Exception:
                                pass
                        # Y axis title
                        if y_axis_cfg.get("title"):
                            try:
                                ax.has_title = True
                                ax.axis_title.text_frame.text = y_axis_cfg["title"]
                                for p in ax.axis_title.text_frame.paragraphs:
                                    for r in p.runs:
                                        r.font.size = tokens["chart"]["axis_title_font_size"]
                                        r.font.name = tokens["font"]["family"]["body"]
                                        r.font.color.rgb = text_color
                                        try:
                                            rPr = r._r.get_or_add_rPr()
                                            _set_east_asia_font(rPr, cn_font)
                                        except Exception:
                                            pass
                            except Exception:
                                pass
                    # number format
                    try:
                        ax.tick_labels.number_format = number_format
                        ax.tick_labels.number_format_is_linked = False
                    except Exception:
                        pass
                except Exception:
                    pass

            _style_axis(c_ax, is_value=False)
            _style_axis(v_ax, is_value=True)
        except Exception:
            pass

    # ---------- Gap width for bar/column (cleaner look) ----------
    try:
        ctype = chart_data["type"]
        if ctype in ("column", "bar", "stacked_column", "stacked_bar"):
            plot = chart.plots[0]
            ser_elem = plot._element
            # gapWidth on the barChart / bar3DChart element
            for child in ser_elem:
                tag = etree.QName(child).localname if hasattr(child, "tag") else ""
                if tag in ("barChart", "bar3DChart"):
                    gw = child.find(qn("c:gapWidth"))
                    if gw is None:
                        gw = etree.SubElement(child, qn("c:gapWidth"))
                    gw.set("val", "80")
    except Exception:
        pass


def _enable_data_labels(series, number_format: str, ctype: str, label_size,
                        text_color, series_color, cn_font: str, tokens: dict):
    """Add/override <c:dLbls> for a series."""
    ser_elem = series._element
    # Remove existing dLbls
    for old in ser_elem.findall(qn("c:dLbls")):
        ser_elem.remove(old)
    dLbls = etree.SubElement(ser_elem, qn("c:dLbls"))
    # show value (and for pie, show category name + percentage)
    if ctype in ("pie", "doughnut"):
        etree.SubElement(dLbls, qn("c:showVal")).set("val", "0")
        etree.SubElement(dLbls, qn("c:showCatName")).set("val", "1")
        etree.SubElement(dLbls, qn("c:showPercent")).set("val", "1")
        etree.SubElement(dLbls, qn("c:showSerName")).set("val", "0")
        etree.SubElement(dLbls, qn("c:showLegendKey")).set("val", "0")
        # Best-fit position outside end
        etree.SubElement(dLbls, qn("c:dLblPos")).set("val", "outEnd")
        numFmt = etree.SubElement(dLbls, qn("c:numFmt"))
        numFmt.set("formatCode", "0%")
        numFmt.set("sourceLinked", "0")
    else:
        etree.SubElement(dLbls, qn("c:showVal")).set("val", "1")
        etree.SubElement(dLbls, qn("c:showCatName")).set("val", "0")
        etree.SubElement(dLbls, qn("c:showSerName")).set("val", "0")
        etree.SubElement(dLbls, qn("c:showLegendKey")).set("val", "0")
        etree.SubElement(dLbls, qn("c:dLblPos")).set("val", "outEnd")
        numFmt = etree.SubElement(dLbls, qn("c:numFmt"))
        numFmt.set("formatCode", number_format)
        numFmt.set("sourceLinked", "0")

    # Text properties
    txPr = etree.SubElement(dLbls, qn("c:txPr"))
    bodyPr = etree.SubElement(txPr, qn("a:bodyPr"))
    lstStyle = etree.SubElement(txPr, qn("a:lstStyle"))
    p = etree.SubElement(txPr, qn("a:p"))
    pPr = etree.SubElement(p, qn("a:pPr"))
    defRPr = etree.SubElement(pPr, qn("a:defRPr"))
    sz_hundredths = int(label_size.pt * 100 if hasattr(label_size, "pt") else (label_size / 12700) * 100)
    defRPr.set("sz", str(sz_hundredths))
    defRPr.set("b", "0")
    # latin font
    latin = etree.SubElement(defRPr, qn("a:latin"))
    latin.set("typeface", tokens["font"]["family"]["body"])
    ea = etree.SubElement(defRPr, qn("a:ea"))
    ea.set("typeface", cn_font)
    solid = etree.SubElement(defRPr, qn("a:solidFill"))
    srgb = etree.SubElement(solid, qn("a:srgbClr"))
    srgb.set("val", "{:02X}{:02X}{:02X}".format(text_color[0], text_color[1], text_color[2]))


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------
def render_chart(slide, chart_data: dict[str, Any], x: int, y: int, cx: int, cy: int,
                 tokens: dict, cn_font: str = "微软雅黑"):
    """Render a native OOXML chart onto ``slide``.

    Parameters
    ----------
    slide : pptx.slide.Slide
        Target slide (returned by ``Renderer.new_slide``).
    chart_data : dict
        Semantic chart spec: type, title, categories, series, show_legend, etc.
    x, y, cx, cy : int
        Bounding box in EMU (calculated by GridLayout).
    tokens : dict
        Fully-resolved theme tokens.
    cn_font : str
        East-Asian font name (default 微软雅黑).
    """
    ctype = chart_data.get("type")
    if ctype not in CHART_TYPE_MAP:
        # English: f"Unsupported chart type: {ctype!r}. Supported: ..."
        raise ValueError(
            f"不支持的图表类型 {ctype!r}。修复建议：请从以下类型中选择：{sorted(CHART_TYPE_MAP)}"
        )
    xl_type = CHART_TYPE_MAP[ctype]

    # Auto-hide legend for single-series column/bar/line
    show_legend = chart_data.get("show_legend")
    if show_legend is None:
        single_series = len(chart_data.get("series", [])) <= 1
        if single_series and ctype in ("column", "bar", "stacked_column", "stacked_bar",
                                       "line", "line_markers", "area", "radar"):
            chart_data = dict(chart_data)  # local copy
            chart_data["show_legend"] = False

    # Build chart data object
    if ctype in _SCATTER_TYPES:
        cd = _build_xy_data(chart_data)
    else:
        cd = _build_category_data(chart_data)

    margin = tokens["chart"]["margin"]
    frame = slide.shapes.add_chart(
        xl_type,
        x + int(margin),
        y + int(margin),
        cx - 2 * int(margin),
        cy - 2 * int(margin),
        cd,
    )
    chart = frame.chart
    _style_chart(chart, chart_data, tokens, xl_type, cn_font)
    return frame
