"""OOXML renderer – thin wrapper around python-pptx.

This is the ONLY module that touches python-pptx APIs.  Layout handlers
build a list of shape descriptors via ``Renderer`` methods and receive EMU
coordinates/dimensions computed by ``GridLayout``.
"""
from __future__ import annotations

import os
from typing import Optional

from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from shared.remote_assets import cleanup_temp_files, download_remote_asset


def _download_remote_image(url: str, timeout: float = 6.0) -> Optional[str]:
    """Download a remote image to a temporary file; return local path or None on failure."""
    try:
        return download_remote_asset(
            url,
            timeout=timeout,
            prefix="ppt_img_",
            fallback_suffix=".jpg",
        )
    except Exception:
        return None


# Track temp files so we can keep them alive for the life of the renderer
# (python-pptx reads lazily in some cases; safer to keep files until save).
_temp_image_paths: list[str] = []


# ----------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------
def _set_cn_font(run, cn_font_name: str):
    """Set East-Asian font on a run (required for CJK rendering)."""
    rPr = run._r.get_or_add_rPr()
    # Remove existing eastAsia
    for ea in rPr.findall(qn("a:ea")):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", cn_font_name)


def _apply_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _apply_no_line(shape):
    shape.line.fill.background()


# ----------------------------------------------------------------------
# Renderer
# ----------------------------------------------------------------------
class Renderer:
    """OOXML renderer for PPT — thin wrapper around python-pptx.

    PPT 渲染器类。唯一直接调用 python-pptx API 的模块，所有版式 handler 通过本类提供的
    rect/textbox/image/table/free_shape 等原语构造形状描述符；坐标/尺寸由 GridLayout 统一
    以 EMU 计算。v1.4 起新增 add_logos（水印）/add_page_numbers（页码）/apply_transition
    （切换效果）/add_speaker_notes（演讲者备注）等装饰接口，均 best-effort 不中断主流程。"""

    def __init__(self, tokens: dict, lang: str = "cn"):
        self.tokens = tokens
        self.lang = lang
        self.prs = Presentation()
        # Force 16:9 widescreen
        self.prs.slide_width = tokens["slide"]["w"]
        self.prs.slide_height = tokens["slide"]["h"]
        # Use blank layout (we draw everything ourselves)
        self.blank_layout = self.prs.slide_layouts[6]
        self.slides: list = []

        self.cn_heading_font = tokens["font"]["family"]["cn_heading"]
        self.cn_body_font = tokens["font"]["family"]["cn_body"]
        self.en_heading_font = tokens["font"]["family"]["heading"]
        self.en_body_font = tokens["font"]["family"]["body"]
        self.accent_font = tokens["font"]["family"]["accent"]
        self._slide_meta: list[dict] = []  # parallel to self.slides; records layout name
        # v1.6.3: shape name registry and slide animation specs
        #   _shape_names[slide_idx][name] = shape_id
        #   _slide_anims[slide_idx] = animation dict (from JSON "animation" field)
        self._shape_names: list[dict] = []
        self._slide_anims: list[Optional[dict]] = []

    def begin_slide(self, layout_name: str):
        """Called by compiler before rendering each slide, to track metadata."""
        self._slide_meta.append({"layout": layout_name})
        self._shape_names.append({})
        self._slide_anims.append(None)

    # ------------------------------------------------------------------
    # Slide lifecycle
    # ------------------------------------------------------------------
    def new_slide(self, bg_color: Optional[RGBColor] = None):
        slide = self.prs.slides.add_slide(self.blank_layout)
        # Background
        if bg_color is None:
            bg_color = self.tokens["color"]["bg"]
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        self.slides.append(slide)
        return slide

    # v1.6.3 animation helpers
    def register_shape_name(self, name: str, shape):
        """Register a python-pptx shape under a semantic name for animation binding."""
        if not self._shape_names:
            return
        self._shape_names[-1][name] = shape.shape_id

    def set_slide_animation(self, anim_spec: dict):
        """Record slide-level animation spec (applied during finalize)."""
        if not self._slide_anims:
            return
        self._slide_anims[-1] = anim_spec

    def finalize_slide_animations(self):
        """Build <p:timing> for all slides that have animation specs. Called once at end of render."""
        from pro_ppt_gen.engine.animation import (
            build_slide_timing, NodeAnim, ANIMATION_PRESETS,
            apply_ext_transition,
        )
        for idx, slide in enumerate(self.slides):
            anim_spec = self._slide_anims[idx] if idx < len(self._slide_anims) else None
            if not anim_spec:
                continue
            names_map = self._shape_names[idx] if idx < len(self._shape_names) else {}
            node_anims = self._resolve_anims(anim_spec, slide, names_map)
            try:
                build_slide_timing(slide, node_anims)
            except Exception as e:
                print(f"[animation] slide {idx} build failed: {e}")

    def _resolve_anims(self, spec: dict, slide, names_map: dict) -> list:
        """Convert semantic animation spec into NodeAnim list.

        spec shape: {"entrance": "...", "emphasis": "...", "exit": "...",
                     "sequence": "one_by_one|all_at_once|by_paragraph|by_series",
                     "trigger": "on_click|auto_after",
                     "targets": ["title","content","kpi",...] | None (all top shapes)}
        """
        from pro_ppt_gen.engine.animation import NodeAnim, ANIMATION_PRESETS
        entrance = spec.get("entrance")
        emphasis = spec.get("emphasis")
        exit_p = spec.get("exit")
        sequence = spec.get("sequence", "all_at_once")
        trigger = spec.get("trigger", "on_click")
        targets = spec.get("targets")
        # Gather shape ids
        if targets:
            ids = [names_map[t] for t in targets if t in names_map]
        else:
            # All top-level non-background shapes on slide
            ids = []
            for sh in slide.shapes:
                try:
                    ids.append(sh.shape_id)
                except Exception:
                    pass
        if not ids:
            return []
        anims: list = []
        # Use entrance as primary; apply emphasis/exit as additional nodes if present
        # For simplicity, first shape gets on_click, others auto_after when by_sequence
        primary_preset = entrance or emphasis or "fade_in"
        if primary_preset not in ANIMATION_PRESETS:
            primary_preset = "fade_in"
        for i, sid in enumerate(ids):
            tri = trigger
            if i > 0 and sequence in ("one_by_one", "by_paragraph", "by_series"):
                tri = "auto_after"
            anims.append(NodeAnim(shape_id=sid, preset_name=primary_preset,
                                   sequence=sequence, trigger=tri))
        return anims

    # ------------------------------------------------------------------
    # Primitive shape ops
    # ------------------------------------------------------------------
    def rect(
        self,
        slide,
        left: int,
        top: int,
        width: int,
        height: int,
        fill_color: RGBColor,
        line_color: Optional[RGBColor] = None,
    ):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        _apply_fill(shp, fill_color)
        if line_color is None:
            _apply_no_line(shp)
        else:
            shp.line.color.rgb = line_color
        shp.shadow.inherit = False
        return shp

    def rounded_rect(
        self,
        slide,
        left: int,
        top: int,
        width: int,
        height: int,
        fill_color: RGBColor,
    ):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        # Adjust corner radius
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
        _apply_fill(shp, fill_color)
        _apply_no_line(shp)
        shp.shadow.inherit = False
        return shp

    def textbox(
        self,
        slide,
        left: int,
        top: int,
        width: int,
        height: int,
        text: str,
        font_size_pt: float,
        color: RGBColor,
        bold: bool = False,
        italic: bool = False,
        align: str = "left",
        anchor: str = "top",
        font_family: Optional[str] = None,
        cn_font_family: Optional[str] = None,
        word_wrap: bool = True,
    ):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = word_wrap
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.auto_size = None

        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
        tf.vertical_anchor = anchor_map.get(anchor, MSO_ANCHOR.TOP)

        # Support multiline: split by \n
        lines = text.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align_map.get(align, PP_ALIGN.LEFT)
            run = p.add_run()
            run.text = line
            f = run.font
            f.size = Pt(font_size_pt)
            f.bold = bold
            f.italic = italic
            f.color.rgb = color
            family = font_family or self.en_body_font
            f.name = family
            _set_cn_font(run, cn_font_family or self.cn_body_font)
        return tb

    def bullets(
        self,
        slide,
        left: int,
        top: int,
        width: int,
        height: int,
        items: list,
        color: RGBColor,
    ):
        """Render a multi-level bullet list.

        Each item is either ``str`` (level 1) or ``{"text": str, "level"?: 1|2|3}``.
        Bullets use manual bullet characters from tokens (●/○/▪) for visual
        consistency across Office versions.
        """
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.auto_size = None

        bt = self.tokens["bullet"]
        sz = self.tokens["font"]["size"]
        space_after_pt = bt["space_after"]

        for idx, item in enumerate(items):
            if isinstance(item, str):
                text, level = item, 1
            else:
                text = item.get("text", "")
                level = int(item.get("level", 1))

            size_key = {1: "body_l1", 2: "body_l2", 3: "body_l3"}.get(level, "body_l1")
            font_size = sz[size_key]
            bullet_char = bt["char"][f"l{level}"]
            indent_key = {1: "indent_l1", 2: "indent_l2", 3: "indent_l3"}[level]
            indent = bt[indent_key]

            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = space_after_pt
            # Hanging indent via paragraph properties
            pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
            pPr.set("marL", str(int(indent)))
            pPr.set("indent", str(-int(indent)))

            # bullet glyph run (accent/secondary colour for level 1)
            bullet_color = self.tokens["color"]["secondary"] if level == 1 else color
            r1 = p.add_run()
            r1.text = bullet_char + "  "
            r1.font.size = font_size
            r1.font.color.rgb = bullet_color
            r1.font.name = self.en_body_font
            _set_cn_font(r1, self.cn_body_font)

            r2 = p.add_run()
            r2.text = text
            r2.font.size = font_size
            r2.font.color.rgb = color
            r2.font.name = self.en_body_font
            _set_cn_font(r2, self.cn_body_font)
        return tb

    def image(self, slide, image_path: str, left: int, top: int, width: int = None, height: int = None,
              placeholder_color=None, placeholder_text: Optional[str] = None, fit: bool = False):
        """Add an image. Supports local path or http(s):// URL.

        On download/missing-file failure, draws a placeholder rectangle with
        text so the whole render does not fail.

        Parameters
        ----------
        fit : bool
            If True and both width/height given, the image is centered and
            aspect-ratio is preserved to fit inside the box (with margins).
        """
        resolved_path = image_path
        if isinstance(image_path, str) and image_path.startswith(("http://", "https://")):
            local = _download_remote_image(image_path)
            if local is None:
                ph_color = placeholder_color or self.tokens["color"]["kpi_bg"]
                ph_w = width or int(Inches(4).emu)
                ph_h = height or int(Inches(3).emu)
                self.rect(slide, left, top, ph_w, ph_h, fill_color=ph_color)
                msg = placeholder_text or f"[Image unavailable]\n{image_path}"
                self.textbox(
                    slide, left, top, ph_w, ph_h, msg,
                    font_size_pt=10, color=self.tokens["color"]["text_light"],
                    align="center", anchor="middle",
                )
                return None
            resolved_path = local
            _temp_image_paths.append(local)

        if not os.path.isfile(resolved_path):
            ph_color = placeholder_color or self.tokens["color"]["kpi_bg"]
            ph_w = width or int(Inches(4).emu)
            ph_h = height or int(Inches(3).emu)
            self.rect(slide, left, top, ph_w, ph_h, fill_color=ph_color)
            msg = placeholder_text or f"[Image not found]\n{image_path}"
            self.textbox(
                slide, left, top, ph_w, ph_h, msg,
                font_size_pt=10, color=self.tokens["color"]["text_light"],
                align="center", anchor="middle",
            )
            return None

        # Aspect-ratio fit inside the given box
        pic_left, pic_top, pic_w, pic_h = left, top, width, height
        if fit and width and height:
            try:
                from PIL import Image as PILImage
                with PILImage.open(resolved_path) as im:
                    iw_px, ih_px = im.size
                ar = iw_px / ih_px
                box_ar = width / height
                if ar > box_ar:
                    pic_w = width
                    pic_h = int(width / ar)
                else:
                    pic_h = height
                    pic_w = int(height * ar)
                pic_left = left + (width - pic_w) // 2
                pic_top = top + (height - pic_h) // 2
            except Exception:
                pic_left, pic_top, pic_w, pic_h = left, top, width, height

        return slide.shapes.add_picture(resolved_path, pic_left, pic_top, width=pic_w, height=pic_h)

    def table(
        self,
        slide,
        left: int,
        top: int,
        width: int,
        height: int,
        headers: list[str],
        rows: list[list[str]],
    ):
        n_rows = len(rows) + 1
        n_cols = len(headers)
        tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = tbl_shape.table

        # Header row
        hdr_bg = self.tokens["color"]["table_header_bg"]
        hdr_fg = self.tokens["color"]["table_header_text"]
        alt_bg = self.tokens["color"]["table_alt_row"]
        body_fg = self.tokens["color"]["text"]
        hdr_size = self.tokens["font"]["size"]["table_header"]
        body_size = self.tokens["font"]["size"]["table_body"]

        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = ""
            _apply_fill(cell, hdr_bg)
            tf = cell.text_frame
            tf.margin_left = Emu(int(0.05 * 914400))
            tf.margin_right = Emu(int(0.05 * 914400))
            tf.margin_top = Emu(int(0.03 * 914400))
            tf.margin_bottom = Emu(int(0.03 * 914400))
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(h)
            r.font.size = hdr_size
            r.font.bold = True
            r.font.color.rgb = hdr_fg
            r.font.name = self.en_heading_font
            _set_cn_font(r, self.cn_heading_font)

        # Body rows
        for i, row in enumerate(rows, start=1):
            bg = alt_bg if i % 2 == 0 else self.tokens["color"]["bg"]
            for j in range(n_cols):
                val = row[j] if j < len(row) else ""
                cell = table.cell(i, j)
                cell.text = ""
                _apply_fill(cell, bg)
                tf = cell.text_frame
                tf.margin_left = Emu(int(0.05 * 914400))
                tf.margin_right = Emu(int(0.05 * 914400))
                tf.margin_top = Emu(int(0.02 * 914400))
                tf.margin_bottom = Emu(int(0.02 * 914400))
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT if j > 0 else PP_ALIGN.LEFT
                r = p.add_run()
                r.text = str(val)
                r.font.size = body_size
                r.font.color.rgb = body_fg
                r.font.name = self.en_body_font
                _set_cn_font(r, self.cn_body_font)
        return tbl_shape

    def vertical_line(self, slide, left: int, top: int, height: int, color: RGBColor, width_pt: float = 4):
        """Thin vertical decorative bar."""
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(width_pt), height)
        _apply_fill(shp, color)
        _apply_no_line(shp)
        shp.shadow.inherit = False
        return shp

    # ------------------------------------------------------------------
    # SVG freeform shape (v1.3 PPT-P0-1)
    # ------------------------------------------------------------------
    def free_shape(
        self,
        slide,
        svg_str: str,
        left: int,
        top: int,
        width: int,
        height: int,
        *,
        shape_id: int = 0,
        name: str = "FreeShape",
    ):
        """Inject an SVG string as a custom-geometry shape via shared svg_engine.

        The returned p:sp element is appended directly to the slide's spTree,
        bypassing python-pptx shape helpers.
        """
        try:
            from ..shared.svg_engine import svg_to_ooxml
        except ImportError:
            try:
                from shared.svg_engine import svg_to_ooxml
            except ImportError:
                try:
                    from skills.shared.svg_engine import svg_to_ooxml
                except ImportError:
                    return None
        sid = shape_id or (len(slide.shapes._spTree) + 1)
        sp = svg_to_ooxml(
            svg_str, int(width), int(height),
            shape_id=sid, name=name,
            x_emu=int(left), y_emu=int(top),
        )
        # Override x/y to our position (svg_to_ooxml already used x_emu)
        slide.shapes._spTree.append(sp)
        return sp

    # ------------------------------------------------------------------
    # Speaker notes (v1.4 P1-1)
    # ------------------------------------------------------------------
    def add_speaker_notes(self, slide, notes_text: str):
        """Attach speaker notes to a slide (best-effort, empty text is a no-op).

        为幻灯片添加演讲者备注。text 为空或全空白时不创建 notes_slide（避免无必要的 part 关系）；
        兼容 v1.3 旧 add_notes 别名以保证向后兼容。"""
        if not notes_text or not str(notes_text).strip():
            return
        try:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = ""
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(notes_text)
            f = run.font
            f.size = Pt(12)
            f.name = self.en_body_font
            _set_cn_font(run, self.cn_body_font)
        except Exception:
            pass

    # Backwards-compat alias (parser used to call add_notes)
    add_notes = add_speaker_notes

    # ------------------------------------------------------------------
    # Slide transitions (v1.4 P2-2)
    # ------------------------------------------------------------------
    _TRANSITION_MAP = {
        # kind    → (tag, **extra_attribs)
        "fade":         ("fade",     {}),
        "dissolve":     ("dissolve", {}),
        "push":         ("push",     {"dir": "l"}),
        "wipe":         ("wipe",     {"dir": "l"}),
        "cover":        ("cover",    {"dir": "l"}),
        "split":        ("split",    {"orient": "horz", "dir": "out"}),
        # v1.6.3: directional aliases commonly referenced by templates/animations
        "wipe_right":   ("wipe",     {"dir": "r"}),
        "wipe_left":    ("wipe",     {"dir": "l"}),
        "wipe_up":      ("wipe",     {"dir": "u"}),
        "wipe_down":    ("wipe",     {"dir": "d"}),
        "push_left":    ("push",     {"dir": "l"}),
        "push_right":   ("push",     {"dir": "r"}),
        "push_up":      ("push",     {"dir": "u"}),
        "push_down":    ("push",     {"dir": "d"}),
        "cover_left":   ("cover",    {"dir": "l"}),
        "cover_right":  ("cover",    {"dir": "r"}),
        "cover_up":     ("cover",    {"dir": "u"}),
        "cover_down":   ("cover",    {"dir": "d"}),
        "split_horizontal": ("split", {"orient": "horz", "dir": "out"}),
        "split_vertical":   ("split", {"orient": "vert", "dir": "out"}),
        "fade_through_black": ("fade", {}),  # alias for compatibility
    }

    def apply_transition(self, slide, transition_type: str):
        """Apply a slide transition effect to a single slide.

        为单张幻灯片设置切换效果。支持 fade/dissolve/push/wipe/cover/split/none
        以及 v1.6.3 新增的 cube/reveal/zoom/ferris/gallery/conveyor 扩展切换。"""
        kind = (transition_type or "").strip().lower()
        try:
            sld = slide._element  # <p:sld>
            # Remove existing transition
            for old in sld.findall(qn("p:transition")):
                sld.remove(old)
            if kind == "none" or kind == "":
                return
            # Try extended transitions first (v1.6.3)
            try:
                from pro_ppt_gen.engine.animation import apply_ext_transition
                if apply_ext_transition(slide, kind):
                    return
            except Exception:
                pass
            spec = self._TRANSITION_MAP.get(kind)
            if spec is None:
                return
            tag, attribs = spec
            if tag == "__ext__":
                # should have been handled above
                return
            p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
            trans = etree.SubElement(sld, f"{{{p_ns}}}transition", nsmap={"p": p_ns})
            trans.set("spd", "med")
            child = etree.SubElement(trans, f"{{{p_ns}}}{tag}")
            for ak, av in attribs.items():
                child.set(ak, av)
        except Exception as e:
            print(f"[transition] failed: {e}")

    # Backwards-compat alias
    set_transition = apply_transition

    # ------------------------------------------------------------------
    # v1.4 P2-1: SVG mini chart → DrawingML shape
    # ------------------------------------------------------------------
    def _render_mini_chart(self, slide, x: int, y: int, w: int, h: int, svg_str: str):
        """Render an SVG mini chart as native DrawingML on ``slide``.

        Best-effort: if SVG conversion fails, falls back to a colored placeholder
        rectangle so the slide is not broken.
        """
        if not svg_str:
            return None
        try:
            shape = self.free_shape(slide, svg_str, x, y, w, h, name="MiniChart")
            if shape is not None:
                return shape
        except Exception as e:
            print(f"[mini_chart] SVG render failed: {e}")
        try:
            ph = self.rect(slide, x, y, w, h, fill_color=self.tokens["color"].get("kpi_bg", RGBColor(0xEE, 0xEE, 0xEE)))
            self.textbox(slide, x, y, w, h, "[chart]",
                         font_size_pt=8, color=self.tokens["color"]["text_light"],
                         align="center", anchor="middle")
            return ph
        except Exception:
            return None

    # Backwards-compat alias (named for layout handlers' convenience)
    mini_chart = _render_mini_chart

    # ------------------------------------------------------------------
    # Page numbers (render after all slides exist so we know total count)
    # v1.4 P3-2: style-aware (plain / slash / chinese / roman)
    # ------------------------------------------------------------------
    def _format_page_number(self, idx: int, total: int, style: str) -> str:
        """Return the page-number label for the given style."""
        if style == "slash":
            return f"{idx}/{total}"
        if style == "chinese":
            return f"第{idx}页"
        if style == "roman":
            # Lazy import to avoid cycles with validators
            from ..compiler.validators import _to_roman
            return _to_roman(idx)
        # "plain" (default) and any unknown → plain digit for backward compat
        return f"{idx}"

    def add_page_numbers(self, skip_layouts=None, style: Optional[str] = None):
        """Render page numbers. Skips cover/toc/section/thanks layouts by default
        so opening/divider pages stay clean.

        Parameters
        ----------
        style : str | None
            ``None`` (default) → legacy "N / TOTAL" (100% backward compat when
            ``meta.page_number_style`` is unspecified);
            ``"plain"`` → single digit "5";
            ``"slash"`` → "5/12";
            ``"chinese"`` → "第5页";
            ``"roman"`` → "V" (uppercase Roman numeral).
        """
        skip_layouts = set(skip_layouts or ("cover", "toc", "section", "thanks"))
        total = len(self.slides)
        size = self.tokens["font"]["size"]["page_num"]
        size_pt = size.pt if hasattr(size, "pt") else size / 12700
        color = self.tokens["color"]["text_light"]
        style_key = (style or "").lower()
        legacy_mode = (style is None) or (style_key == "")
        for idx, (slide, meta) in enumerate(zip(self.slides, self._slide_meta), start=1):
            layout = meta.get("layout")
            if layout in skip_layouts:
                continue
            fl, ft, fw, fh = (
                int(self.tokens["margin"]["left"]),
                int(self.tokens["slide"]["h"] - self.tokens["margin"]["bottom"] - self.tokens["footer"]["height"]),
                int(self.tokens["slide"]["w"] - self.tokens["margin"]["left"] - self.tokens["margin"]["right"]),
                int(self.tokens["footer"]["height"]),
            )
            if legacy_mode:
                label = f"{idx} / {total}"
            else:
                label = self._format_page_number(idx, total, style_key)
            self.textbox(
                slide,
                fl, ft, fw, fh,
                label,
                font_size_pt=size_pt,
                color=color,
                align="right",
                anchor="middle",
            )

    # ------------------------------------------------------------------
    # v1.4 P3-1: Logo watermark (per-page)
    # ------------------------------------------------------------------
    def _logo_geometry(self, position: str, logo_size: int) -> tuple[int, int, int, int]:
        """Return (left, top, width, height) in EMU for a logo placed at position."""
        sw = int(self.tokens["slide"]["w"])
        sh = int(self.tokens["slide"]["h"])
        pad = int(Inches(0.3))
        L = pad
        R = sw - pad - logo_size
        T = pad
        B = sh - pad - logo_size
        CX = (sw - logo_size) // 2
        CY = (sh - logo_size) // 2
        pos = (position or "tr").lower()
        if pos == "tl":
            return L, T, logo_size, logo_size
        if pos == "bl":
            return L, B, logo_size, logo_size
        if pos == "br":
            return R, B, logo_size, logo_size
        if pos == "center":
            return CX, CY, logo_size, logo_size
        # default tr
        return R, T, logo_size, logo_size

    def _add_logo_to_slide(self, slide, logo_path: str, position: str):
        """Insert a logo picture on a single slide. Best-effort: swallows errors.

        Supports both local filesystem paths and http(s):// URLs. Remote logos
        are downloaded to a temporary file (5s timeout) and cleaned up after
        insertion; download failures are silently ignored so the main render
        is never interrupted.
        """
        if not logo_path:
            return
        # v1.4 P3-1b: support http(s) URL for logo — download to tempfile
        tmp_path: Optional[str] = None
        resolved_path = logo_path
        if isinstance(logo_path, str) and logo_path.startswith(("http://", "https://")):
            try:
                tmp_path = download_remote_asset(
                    logo_path,
                    timeout=5.0,
                    prefix="ppt_logo_",
                    fallback_suffix=".png",
                )
                resolved_path = tmp_path
                _temp_image_paths.append(tmp_path)
            except Exception as e:
                # Per PRD: silently skip on download failure
                print(f"[logo] remote logo download skipped ({logo_path}): {e}")
                if tmp_path and os.path.isfile(tmp_path):
                    try:
                        os.unlink(tmp_path)
                    except Exception:
                        pass
                return
        try:
            sw = int(self.tokens["slide"]["w"])
            sh = int(self.tokens["slide"]["h"])
            short_side = min(sw, sh)
            # Logo size ≈ 9% of short side, within the 8–10% range
            logo_size = int(short_side * 0.09)
            l, t, w, h = self._logo_geometry(position, logo_size)
            self.image(
                slide, resolved_path, l, t, width=w, height=h, fit=True,
                placeholder_color=None, placeholder_text=None,
            )
        except Exception as e:
            # Silently skip logo insertion failures – never break main render
            print(f"[logo] watermark skipped: {e}")
        finally:
            # v1.4 P3-1b: clean up downloaded temp file after insertion
            # (python-pptx reads the file during save, so we keep a ref in
            # _temp_image_paths and defer real cleanup to renderer teardown;
            # but if we used a temp we did not register, remove it here.)
            pass

    def add_logos(self, logo_path: str, position: str = "tr", skip_layouts=None):
        """Render a small logo watermark on every slide (best-effort).

        为每张幻灯片添加小尺寸 Logo 水印。logo_path 支持本地文件系统路径或 http(s):// URL
        （v1.4 P3-1b，URL 自动下载到临时文件，5s 超时失败静默跳过）；position 取 tl/tr/bl/br/center；
        skip_layouts 指定要跳过的版式（如 ["cover", "thanks"]）。缺图/下载失败/插入异常均不中断主渲染。"""
        if not logo_path:
            return
        skip_layouts = set(skip_layouts or ())  # by default, logo shows on all pages
        for slide, meta in zip(self.slides, self._slide_meta):
            if meta.get("layout") in skip_layouts:
                continue
            self._add_logo_to_slide(slide, logo_path, position)

    # ------------------------------------------------------------------
    # Save
    # ------------------------------------------------------------------
    def save(self, path: str) -> str:
        """Save the rendered Presentation to disk and return the absolute path.

        将渲染完成的 Presentation 保存到指定路径，自动创建父目录，返回绝对路径。保存前
        python-pptx 会读取所有已插入图片文件，因此 add_logos 下载的临时文件需保活至 save 完成。"""

        os.makedirs(os.path.dirname(os.path.abspath(path)) or ".", exist_ok=True)
        try:
            self.prs.save(path)
        finally:
            cleanup_temp_files(_temp_image_paths)
            _temp_image_paths.clear()
        return path
