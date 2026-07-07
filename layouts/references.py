"""References layout — 参考文献页。

v1.4 P1-4b 新增。每页显示标题 + 多条参考文献条目，方括号编号 [N]，小字号（9pt）。
条目过多时自动追加「参考文献（续）」页（进阶分页），单页装不下时先截断并 warning。
"""
from __future__ import annotations

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from ._helpers import emu

# Citation engine shared with DOCX
def _load_cite_format():
    """Lazy-load shared citation formatter."""
    try:
        from ...shared.citation import format_reference  # type: ignore
        return format_reference
    except Exception:
        pass
    try:
        from shared.citation import format_reference  # type: ignore
        return format_reference
    except Exception:
        try:
            # Fallback sys.path to skills root
            import os as _os, sys as _sys
            _here = _os.path.dirname(_os.path.abspath(__file__))
            _skills = _os.path.normpath(_os.path.join(_here, "..", "..", ".."))
            if _skills not in _sys.path:
                _sys.path.insert(0, _skills)
            from shared.citation import format_reference  # type: ignore
            return format_reference
        except Exception:
            def _fallback(item, style="apa", index=None):
                if isinstance(item, dict):
                    return item.get("title", str(item))
                return str(item)
            return _fallback


_CITE = None


def _cite(item, style, index=None):
    global _CITE
    if _CITE is None:
        _CITE = _load_cite_format()
    return _CITE(item, style=style, index=index)


# 单页大约能容纳的条目数（按 9pt/1.15 行距估算；Inches(5) 内容区 ÷ ~0.22 英寸/行）
_ITEMS_PER_PAGE = 12


def _render_page(slide, renderer, grid, tokens, title, lines, is_continue, start_idx):
    """渲染单页参考文献。"""
    # 标题（col 1-11，顶部 0.5")
    _, t_top, t_w, t_h = grid.col_box(1, 11, emu(Inches(0.4)), emu(Inches(0.7)))
    title_text = title + ("（续）" if is_continue else "")
    if renderer.lang != "cn":
        title_text = title + (" (cont.)" if is_continue else "")
    title_size = tokens["font"]["size"].get("slide_title", Pt(28))
    # slide_title 是 EMU(Pt(...))，textbox 需要 pt 值
    title_pt = title_size / 12700 if title_size > 1000 else title_size
    renderer.textbox(
        slide,
        emu(grid.col_x(1)), emu(t_top), emu(t_w), emu(t_h),
        title_text,
        font_size_pt=title_pt,
        color=tokens["color"].get("heading") or tokens["color"].get("primary") or RGBColor(0x1F, 0x38, 0x64),
        bold=True,
        align="left",
        anchor="top",
        font_family=renderer.en_heading_font,
        cn_font_family=renderer.cn_heading_font,
    )

    # 装饰横线（用扁矩形模拟）
    line_y = t_top + t_h + emu(Inches(0.05))
    line_w = emu(grid.col_x(11) - grid.col_x(1))
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        emu(grid.col_x(1)), emu(line_y), line_w, Pt(2))
    # Lazy import to avoid circular
    from pro_ppt_gen.engine.renderer import _apply_fill, _apply_no_line
    _apply_fill(shp, tokens["color"]["accent"])
    _apply_no_line(shp)
    shp.shadow.inherit = False

    # 内容区
    body_top = line_y + emu(Inches(0.2))
    body_bottom = emu(Inches(7.0))
    body_h = body_bottom - body_top
    body_l = grid.col_x(1)
    body_w = grid.col_x(11) - body_l

    # 拼文本，每行 "[N] xxx"
    numbered = True  # PPT references 默认方括号编号
    text_lines = []
    for i, line in enumerate(lines):
        idx = start_idx + i
        prefix = f"[{idx}] " if numbered else ""
        text_lines.append(prefix + line)
    body_text = "\n".join(text_lines)

    # 条目字体：中英混合 9pt（en_body + cn_body）
    # textbox 内逐段渲染支持 multiline
    tb = slide.shapes.add_textbox(emu(body_l), emu(body_top), emu(body_w), emu(body_h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.auto_size = None

    # 设置行距 1.15
    for i, ln in enumerate(text_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(9)
        c = tokens["color"]["text"]
        run.font.color.rgb = c
        run.font.name = renderer.en_body_font
        # East Asian font
        rPr = run._r.get_or_add_rPr()
        for ea in rPr.findall(qn("a:ea")):
            rPr.remove(ea)
        ea_el = etree.SubElement(rPr, qn("a:ea"))
        ea_el.set("typeface", renderer.cn_body_font)


def render(slide_dict: dict, renderer, grid, tokens):
    """Render references slide(s).

    渲染参考文献页：标题 + 编号条目列表，9pt 字号；超过单页容量时自动追加续页。
    """
    items = slide_dict.get("items", []) or []
    title = slide_dict.get("title") or ("参考文献" if renderer.lang == "cn" else "References")
    style = (slide_dict.get("citation_style") or "gb7714").lower()
    style = {"harvard": "apa", "gbt7714": "gb7714"}.get(style, style)
    if style not in {"apa", "gb7714", "mla", "ieee"}:
        style = "gb7714"

    numbered = style in ("gb7714", "ieee")

    # 格式化所有条目
    formatted = []
    for i, ref in enumerate(items, 1):
        idx = i if numbered else None
        try:
            txt = _cite(ref, style=style, index=idx)
        except Exception:
            txt = str(ref.get("title", ref) if isinstance(ref, dict) else ref)
        formatted.append(txt)

    # 分页（简单估算：按 _ITEMS_PER_PAGE 每页；进阶：测文本高后自动分页，本版本先按估算）
    per_page = _ITEMS_PER_PAGE
    chunks = [formatted[i:i+per_page] for i in range(0, len(formatted), per_page)]
    # 如果多于 1 页，warning
    if len(chunks) > 1:
        print(f"[references] {len(formatted)} items -> {len(chunks)} pages (auto-paginated)")

    for page_no, chunk in enumerate(chunks):
        if page_no == 0:
            slide = renderer.new_slide(bg_color=tokens["color"]["bg"])
        else:
            slide = renderer.new_slide(bg_color=tokens["color"]["bg"])
        start_idx = page_no * per_page + (1 if numbered else 0)
        is_continue = page_no > 0
        _render_page(slide, renderer, grid, tokens, title, chunk, is_continue, start_idx)
