"""Public API for the PaperJSX-based pro-ppt-gen engine.

Usage::

    from skills.pro_ppt_gen import ppt_jsx

    ol      = ppt_jsx.outline(content)
    c2      = ppt_jsx.auto_layout(content)            # v1.2: infer layouts
    report  = ppt_jsx.estimate_length(c2)
    qc      = ppt_jsx.quality_check(c2, theme="academic")  # v1.3
    taste   = ppt_jsx.taste_check(c2, theme="academic")    # v1.5
    story   = ppt_jsx.story_check(c2, method="pyramid")     # v1.3
    prompts = ppt_jsx.collect_image_prompts(c2)             # v1.3
    out     = ppt_jsx.generate(c2, "deck.pptx", theme="tech", lang="cn")

    # Modify existing decks (v1.2):
    ppt_jsx.update_slide("out.pptx", 3, {"layout": "content", "title": "...", "bullets": [...]})
    ppt_jsx.add_slide("out.pptx", {"layout": "thanks"}, position=-1)
    ppt_jsx.delete_slide("out.pptx", 5)

@since v1.3.0
- quality_check (PPT-P0-3 RUBRIC 自检)
- story_check (PPT-P0-4 叙事方法论)
- collect_image_prompts (PPT-P0-2 AI 配图通道)
- assertive_title helper
- free_shape / background_shape / decorations / image_prompt 语义字段
"""
from __future__ import annotations

import copy
import os
import posixpath
import re
import tempfile
import zipfile
from typing import Optional, Union

from lxml import etree

from . import __version__ as _PKG_VERSION
from .tokens.design_tokens import get_tokens
from .tokens.themes import ThemeFactory
from .engine.layout import GridLayout, bullets_height_estimate
from .engine.renderer import Renderer
from shared.import_helper import import_shared
from .compiler.parser import compile_slides
from .compiler.validators import validate_slides, ValidationError, SCHEMA


__all__ = [
    "generate", "outline", "estimate_length", "auto_layout",
    "update_slide", "add_slide", "delete_slide",
    "ThemeFactory", "quality_check", "taste_check", "story_check",
    "visual_redesign_check", "build_visual_redesign_prompt", "apply_visual_redesign_guidance",
    "collect_image_prompts", "assertive_title",
]


# ----------------------------------------------------------------------
# Title keys for outline display (per layout)
# ----------------------------------------------------------------------
_TITLE_KEYS = {
    "cover": "title",
    "toc": "title",
    "section": "title",
    "content": "title",
    "two_column": "title",
    "image_text": "title",
    "full_image": "title",
    "table": "title",
    "chart": "title",
    "kpi": "title",
    "quote": "text",
    "summary": "title",
    "thanks": "title",
    "timeline": "title",
}


def _resolve_tokens(theme, lang: str) -> dict:
    tokens = get_tokens(theme)
    tokens["_lang"] = lang
    return tokens


# ----------------------------------------------------------------------
# v1.6.3: deep-merge helper for template theme_overrides / scanner tokens.
# Recursively merges ``overrides`` dict into ``tokens`` in-place:
#   - if both sides are dict → recurse;
#   - else tokens[k] = overrides[k] (scalar/list replaces).
# Skips keys whose name starts with ``_`` (private metadata) so internal
# flags (``_lang``, ``_warnings``, ...) are not overwritten by templates.
# Auto-converts "#RRGGBB" hex strings under the "color" key to RGBColor so
# template data files can stay plain-data without python-pptx imports.
# ----------------------------------------------------------------------
def _coerce_color_value(v):
    from pptx.dml.color import RGBColor
    if isinstance(v, str) and len(v) == 7 and v.startswith("#"):
        try:
            return RGBColor(int(v[1:3], 16), int(v[3:5], 16), int(v[5:7], 16))
        except ValueError:
            return v
    if isinstance(v, list):
        return [_coerce_color_value(x) for x in v]
    return v


def _deep_merge_tokens(tokens: dict, overrides: Optional[dict], *, _in_color: bool = False) -> None:
    if not overrides:
        return
    in_color = _in_color
    for k, v in overrides.items():
        if isinstance(k, str) and k.startswith("_"):
            continue
        if k == "color" and isinstance(v, dict) and isinstance(tokens.get("color"), dict):
            _deep_merge_tokens(tokens["color"], v, _in_color=True)
            continue
        if isinstance(v, dict) and isinstance(tokens.get(k), dict):
            _deep_merge_tokens(tokens[k], v, _in_color=in_color)
        else:
            tokens[k] = _coerce_color_value(v) if in_color else v


# ----------------------------------------------------------------------
# v1.4 P2-5: default output filename helper
# ----------------------------------------------------------------------
_BAD_FN_CHARS = re.compile(r'[\\/:*?"<>|\s]+')


def _default_output_path(content: dict, ext: str = "pptx") -> str:
    """根据 meta.title 生成默认输出文件名：{title}_v{version}.{ext}。

    - title_sanitized: meta.title (缺失用 "untitled")，去特殊字符/空白用下划线替换，截断40字符
    - version: 包 __version__ 取主版本号（如 1.4.0 → v1.4；1.4 → v1.4）
    """
    meta = (content or {}).get("meta") or {}
    title = meta.get("title") or ""
    # PPT 顶层字段也兼容（有些调用把 slides 当顶层，title 放在 content.title / cover.title）
    if not title and isinstance(content, dict):
        for s in content.get("slides", []) or []:
            if s.get("layout") == "cover" and s.get("title"):
                title = s["title"]; break
        if not title:
            title = content.get("title") or "untitled"
    safe = _BAD_FN_CHARS.sub("_", title).strip("_")
    if not safe:
        safe = "untitled"
    if len(safe) > 40:
        safe = safe[:40].rstrip("_")
    ver = _PKG_VERSION
    # 只取前两段版本号（如 "1.4.0" → "1.4"）
    parts = ver.split(".")
    short_ver = ".".join(parts[:2]) if len(parts) >= 2 else ver
    return f"{safe}_v{short_ver}.{ext}"


# ----------------------------------------------------------------------
# auto_layout – infer layout from semantic fields (v1.2)
# ----------------------------------------------------------------------
def auto_layout(content: dict, in_place: bool = False) -> dict:
    """Infer ``layout`` for each slide based on present semantic fields."""
    if not isinstance(content, dict) or "slides" not in content:
        raise ValidationError("content 必须是包含 'slides' 列表的 dict。修复建议：请传入 {\"slides\": [...]} 结构的对象")
    out = content if in_place else copy.deepcopy(content)

    for slide in out["slides"]:
        if slide.get("layout") and slide["layout"] in SCHEMA:
            continue

        keys = set(slide.keys())
        has_img = ("image_path" in keys) or ("image_prompt" in keys)
        has_chart = "chart" in keys
        has_bullets = bool(slide.get("bullets"))
        has_table = "headers" in keys and "rows" in keys
        has_lr = "left" in keys and "right" in keys
        has_items = "items" in keys
        has_events = "events" in keys
        has_number_title = "number" in keys and "title" in keys
        has_quote_text = "text" in keys and "title" not in keys
        # Cover hints: title + meta fields (subtitle/author/...) without body content
        is_cover = (
            "title" in keys
            and not has_bullets and not has_chart and not has_table
            and not has_lr and not has_items and not has_events
            and not has_quote_text and not has_number_title
            and any(k in keys for k in ("subtitle", "author", "institution", "date"))
        )
        # Thanks hints: explicit thanks-style title only (don't guess)
        is_thanks = slide.get("title") in ("感谢聆听", "谢谢", "Thanks", "Q & A", "谢谢观看", "Thank You")

        if is_cover:
            slide["layout"] = "cover"
        elif is_thanks:
            slide["layout"] = "thanks"
        elif has_events:
            slide["layout"] = "timeline"
        elif has_img and not has_bullets and not has_chart and not has_table:
            slide["layout"] = "full_image"
        elif has_img and has_bullets:
            slide["layout"] = "image_text"
        elif has_chart and not has_bullets:
            slide["layout"] = "chart"
        elif has_table:
            slide["layout"] = "table"
        elif has_lr:
            slide["layout"] = "two_column"
        elif has_items:
            items = slide.get("items") or []
            if items and isinstance(items[0], dict) and ("value" in items[0] or "label" in items[0]):
                slide["layout"] = "kpi"
            else:
                slide["layout"] = "toc"
                slide.setdefault("title", "目  录")
        elif has_number_title and len(keys - {"number", "title", "layout", "notes"}) == 0:
            slide["layout"] = "section"
        elif has_quote_text and not has_bullets:
            slide["layout"] = "quote"
        else:
            slide["layout"] = "content"
            if not has_bullets and not has_chart:
                slide.setdefault("bullets", [])
    return out


# ----------------------------------------------------------------------
# outline – lightweight preview
# ----------------------------------------------------------------------
def outline(content: dict) -> dict:
    """Return a lightweight outline (slide index / layout / title-like field)."""
    if not isinstance(content, dict) or "slides" not in content:
        raise ValidationError("content 必须是包含 'slides' 列表的 dict。修复建议：请传入 {\"slides\": [...]} 结构的对象")
    pages = []
    for i, s in enumerate(content["slides"]):
        layout = s.get("layout", "?")
        key = _TITLE_KEYS.get(layout, "title")
        title_val = s.get(key) or s.get("title") or ""
        if layout == "quote" and title_val and len(title_val) > 40:
            title_val = title_val[:38] + "…"
        pages.append({"index": i + 1, "layout": layout, "title": title_val})
    return {"total_slides": len(pages), "pages": pages}


# ----------------------------------------------------------------------
# estimate_length – overflow + density warnings (v1.2 adds density)
# ----------------------------------------------------------------------
def _density_score(slide: dict, tokens: dict, grid: GridLayout) -> str:
    """Return 'low'/'medium'/'high' based on content volume."""
    layout = slide.get("layout")
    body_h = grid.body_height()
    body_w = grid.content_w
    fs = tokens["font"]["size"]["body_l1"]
    fs_pt = fs.pt if hasattr(fs, "pt") else fs / 12700

    score = 0  # rough "lines of text equivalent"

    bullets = slide.get("bullets") or []
    if bullets:
        h = bullets_height_estimate(bullets, fs_pt, body_w)
        score += h / int(fs_pt * 12700 * 1.4)

    if layout == "two_column":
        for side in ("left", "right"):
            col = slide.get(side) or {}
            b = col.get("bullets", [])
            if b:
                cw = grid.col_w(5)
                h = bullets_height_estimate(b, fs_pt, cw)
                score += h / int(fs_pt * 12700 * 1.4)

    if slide.get("chart"):
        score += 8

    rows = slide.get("rows") or []
    if rows:
        score += len(rows) * 1.2

    events = slide.get("events") or []
    if events:
        score += len(events) * 2

    items = slide.get("items") or []
    if layout == "kpi":
        score += len(items) * 2
    elif layout == "toc":
        score += len(items) * 1.2

    # Title/meta don't count heavily
    if layout == "quote":
        text = slide.get("text", "")
        score = max(score, len(text) / 30)

    if score < 3:
        return "low"
    if score > 14:
        return "high"
    return "medium"


def estimate_length(content: dict) -> dict:
    """Estimate rendered length and flag slides that may overflow or be sparse.

    Returns ``total_slides``, per-slide ``warnings``, ``density`` (low/medium/high),
    and lists of ``overflow_slides`` / ``dense_slides`` / ``sparse_slides``.
    """
    warnings_raw = validate_slides(content)
    tokens = get_tokens("academic")
    grid = GridLayout(tokens)
    pages = []
    overflow_idx = []
    dense_idx = []
    sparse_idx = []
    for i, s in enumerate(content["slides"]):
        page_warnings = []
        layout = s.get("layout")
        body_h = grid.body_height()
        body_w = grid.content_w

        if layout in ("content", "summary"):
            bullets = s.get("bullets", [])
            fs = tokens["font"]["size"]["body_l1"]
            fs_pt = fs.pt if hasattr(fs, "pt") else fs / 12700
            h = bullets_height_estimate(bullets, fs_pt, body_w)
            if h > body_h:
                page_warnings.append(
                    f"bullets likely overflow; reduce bullet count or text length"
                )

        elif layout == "two_column":
            fs = tokens["font"]["size"]["body_l1"]
            fs_pt = fs.pt if hasattr(fs, "pt") else fs / 12700
            col_w = grid.col_w(5)
            for side in ("left", "right"):
                bullets = (s.get(side) or {}).get("bullets", [])
                h = bullets_height_estimate(bullets, fs_pt, col_w)
                if h > body_h:
                    page_warnings.append(f"{side} column bullets may overflow")

        elif layout == "timeline":
            events = s.get("events", [])
            if len(events) > 6:
                page_warnings.append(f"{len(events)} events may crowd timeline; recommend ≤6")

        elif layout == "table":
            rows = s.get("rows", [])
            if len(rows) > 10:
                page_warnings.append(f"{len(rows)} rows may crowd the page")

        elif layout == "chart":
            bullets = s.get("bullets", [])
            if bullets:
                fs = tokens["font"]["size"]["body_l2"]
                fs_pt = fs.pt if hasattr(fs, "pt") else fs / 12700
                h = bullets_height_estimate(bullets, fs_pt, body_w)
                cap = int(body_h * 0.35)
                if h > cap:
                    page_warnings.append("chart insight bullets may overflow; reduce to ≤3 short bullets")

        if layout == "content" and s.get("chart") and s.get("bullets"):
            if len(s.get("bullets", [])) > 3:
                page_warnings.append("bullets+chart coexist; recommend ≤3 bullets")

        for w in warnings_raw:
            if w.slide_idx == i:
                page_warnings.append(w.message)

        density = _density_score(s, tokens, grid)
        if density == "high" and layout not in ("cover", "section", "thanks"):
            page_warnings.append("high information density – consider splitting or simplifying")
            dense_idx.append(i + 1)
        elif density == "low" and layout not in ("cover", "section", "thanks", "quote", "full_image"):
            page_warnings.append("low information density – consider merging or adding content")
            sparse_idx.append(i + 1)

        if page_warnings:
            overflow_idx.append(i + 1)

        # v1.3 PPT-P0-2: image_prompt present but no image_path → "待配图" flag
        has_prompt = s.get("image_prompt") and not s.get("image_path")
        for side in ("left", "right"):
            col = s.get(side)
            if isinstance(col, dict) and col.get("image_prompt") and not col.get("image_path"):
                has_prompt = True
        pending_images = "pending" if has_prompt else "ready"
        if has_prompt:
            page_warnings.append("image_prompt 待生图；调用 collect_image_prompts() 批量处理")

        pages.append({
            "index": i + 1,
            "layout": layout,
            "warnings": page_warnings,
            "density": density,
            "image_status": pending_images,
        })

    return {
        "total_slides": len(pages),
        "pages": pages,
        "overflow_slides": overflow_idx,
        "dense_slides": dense_idx,
        "sparse_slides": sparse_idx,
    }


# ----------------------------------------------------------------------
# generate – compile + render
# ----------------------------------------------------------------------
def generate(
    content: dict,
    output_path: Optional[str] = None,
    theme: Union[str, dict] = "academic",
    lang: str = "cn",
    template_name: Optional[str] = None,
    template_path: Optional[str] = None,
    auto_taste_match: bool = False,
) -> str:
    """编译语义 ``content`` 并渲染 .pptx 到 ``output_path``。

    Compile semantic ``content`` and render a .pptx at ``output_path``.

    Parameters
    ----------
    content : dict
        Semantic JSON with a top-level ``slides`` list (and optional ``meta``).
    output_path : str | None
        Destination .pptx path. 若为 None（v1.4 默认），按 ``{title}_v{version}.pptx`` 规则自动命名。
        Parent dirs created if missing.
    theme : str | dict
        Theme name (e.g. ``"academic"``, ``"tech"``, ``"dark"``) or a custom
        theme dict / ``{"name": "...", "overrides": {...}}``.
    lang : str
        ``"cn"`` picks 微软雅黑 for CJK glyphs; ``"en"`` uses Latin fonts.
    template_name : str | None
        v1.6.3 新增：内置预设模板名（见 pro_ppt_gen.templates.list_templates()）。
        指定后会覆盖 theme 为模板对应的 base_theme + theme_overrides + 默认 transition/animation。
    template_path : str | None
        v1.6.3 新增：本地 .pptx/.potx/.dpt 模板路径，自动提取色板/字体/背景注入 tokens。
    auto_taste_match : bool
        v1.6.3 新增：True 时若模板与内容调性匹配度低（taste_check<60）给出警告与推荐。
    """
    if output_path is None:
        output_path = _default_output_path(content, "pptx")

    # v1.6.3: resolve template_name → base_theme + overrides
    template_obj = None
    if template_name:
        try:
            from .templates import get_template as _gt
            template_obj = _gt(template_name)
            if template_obj is not None:
                if isinstance(theme, str):
                    theme = template_obj.base_theme
        except Exception as e:
            print(f"[template] template_name={template_name} resolve failed: {e}")

    tokens = _resolve_tokens(theme, lang)

    # v1.6.3: apply template theme_overrides
    if template_obj is not None:
        _deep_merge_tokens(tokens, template_obj.theme_overrides)
        # apply default transition/animation per slide if slide didn't specify
        if template_obj.default_transition or template_obj.default_animation:
            for sl in content.get("slides", []):
                if not isinstance(sl, dict):
                    continue
                if template_obj.default_transition and not sl.get("transition"):
                    sl["transition"] = template_obj.default_transition
                if template_obj.default_animation and not sl.get("animation"):
                    sl["animation"] = template_obj.default_animation

    # v1.6.3: apply local template_path theme extraction
    if template_path:
        try:
            from shared.template_scanner import extract_template_theme, apply_extracted_theme
            extracted = extract_template_theme(template_path)
            if extracted.confidence >= 30:
                from shared.themes import THEME_PALETTES
                base_name = template_obj.base_theme if template_obj else (theme if isinstance(theme, str) else "academic")
                merged = apply_extracted_theme(extracted, base_theme=base_name)
                # merge merged tokens into tokens (color/font keys only)
                # apply_extracted_theme 返回的是扁平格式（primary/bg/text/...），需要映射到 tokens.color 子树
                _color_keys = {"primary", "secondary", "accent", "bg", "text",
                               "text_secondary", "bg_secondary", "link",
                               "title_bar_bg", "title_text"}
                _font_family_keys = {"heading_font_latin", "heading_font_ea",
                                     "body_font_latin", "body_font_ea"}
                for k, v in merged.items():
                    if k.startswith("_"):
                        continue
                    if k in _color_keys:
                        tokens.setdefault("color", {})[k] = _coerce_color_value(v)
                    elif k in _font_family_keys:
                        tokens.setdefault("font", {}).setdefault("family", {})[k] = v
                for w in (merged.get("_warnings") or []):
                    print(f"[template] {w}")
            else:
                for w in (extracted.warnings or []):
                    print(f"[template] warning: {w}")
        except Exception as e:
            print(f"[template] extract_template_theme failed: {e}")

    # v1.6.3: auto_taste_match check (lightweight template-match score)
    _match_report = None
    if auto_taste_match:
        try:
            _match_report = check_template_match(
                tokens, template_name=template_name,
                extracted=extracted if template_path else None,
            )
            if not _match_report.get("passed", True):
                for w in _match_report.get("warnings", []):
                    print(f"[template-match] {w}")
        except Exception as e:
            print(f"[template-match] check failed: {e}")

    grid = GridLayout(tokens)
    renderer = Renderer(tokens, lang=lang)

    compile_slides(content, renderer, grid, tokens)

    # v1.4 P1-4b: 若 meta.references 存在（doc 级参考文献），自动追加 references 页
    meta = content.get("meta") or {}
    refs = meta.get("references") if isinstance(meta, dict) else None
    if refs and isinstance(refs, list) and refs:
        # 避免已有 references layout 时重复追加
        already_has = any(
            isinstance(s, dict) and s.get("layout") == "references"
            for s in content.get("slides", [])
        )
        if not already_has:
            ref_slide = {
                "layout": "references",
                "title": meta.get("references_title") or ("参考文献" if lang in ("cn", "zh") else "References"),
                "citation_style": meta.get("citation_style", "gb7714"),
                "items": refs,
            }
            content["slides"].append(ref_slide)
            # re-invoke compile for that single slide via mini content
            try:
                mini = {"meta": meta, "slides": [ref_slide]}
                # Use parser to handle this single slide directly
                from .compiler.parser import _get_handler
                handler = _get_handler("references")
                if hasattr(renderer, "begin_slide"):
                    renderer.begin_slide("references")
                handler(ref_slide, renderer, grid, tokens)
                if hasattr(renderer, "add_speaker_notes"):
                    try:
                        renderer.add_speaker_notes(renderer.slides[-1], "")
                    except Exception:
                        pass
            except Exception as e:
                print(f"[references] auto-append skipped: {e}")

    # v1.4 P3-1: logo watermark (best-effort, failures swallowed by renderer)
    logo_path = meta.get("logo") if isinstance(meta, dict) else None
    logo_position = meta.get("logo_position", "tr") if isinstance(meta, dict) else "tr"
    if logo_path:
        try:
            renderer.add_logos(logo_path, position=logo_position)
        except Exception as e:
            print(f"[logo] add_logos skipped: {e}")

    # v1.4 P3-2: page number style (pass None for legacy compat when unspecified)
    page_style = meta.get("page_number_style") if isinstance(meta, dict) else None
    # If caller explicitly passed plain/slash/chinese/roman, use it; else None → legacy "N / TOTAL"
    try:
        renderer.add_page_numbers(style=page_style)
    except Exception as e:
        print(f"[page_num] skipped: {e}")

    out = os.path.abspath(output_path)
    os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
    renderer.save(out)
    return out


def generate_from_markdown(
    md_text: str,
    output_path: Optional[str] = None,
    theme: Union[str, dict] = "academic",
    lang: str = "cn",
    template_name: Optional[str] = None,
    template_path: Optional[str] = None,
    auto_taste_match: bool = False,
    **kwargs,
) -> str:
    """从 Markdown 文本直接生成 .pptx（P1-2, v1.4；v1.6.3 支持 template_* 参数）。

    Generate a .pptx directly from a Markdown string.

    Args:
        md_text: Markdown 源文本（# → 封面，## → 章节页，-/* → bullets，--- → 分页）。
        output_path: 输出 .pptx 路径，None 时按 ``{title}_v{version}.pptx`` 自动命名。
        theme: 主题名或自定义 tokens dict。
        lang: 语言 ``"cn"`` | ``"en"``。
        template_name/template_path/auto_taste_match: v1.6.3 新参数，透传给 generate()。
        **kwargs: 作为 meta 透传（如 author/date/subtitle）。

    Returns:
        输出文件绝对路径。
    """
    from .compiler.markdown_to_slides import markdown_to_slides
    content = markdown_to_slides(md_text, theme=theme if isinstance(theme, str) else "academic", **kwargs)
    if output_path is None:
        output_path = _default_output_path(content, "pptx")
    return generate(content, output_path=output_path, theme=theme, lang=lang,
                    template_name=template_name, template_path=template_path,
                    auto_taste_match=auto_taste_match)


# ----------------------------------------------------------------------
# Slide modification APIs (v1.2)
# ----------------------------------------------------------------------
def list_themes() -> list:
    """返回所有可用主题的简要信息。"""
    from .tokens.themes import list_themes as _lt
    return [{"name": n} for n in _lt()]


def _render_single_slide(
    slide_spec: dict,
    output_path: str,
    theme: Union[str, dict] = "academic",
    lang: str = "cn",
) -> str:
    """Render a single-slide temp deck containing the given slide spec."""
    tokens = _resolve_tokens(theme, lang)
    grid = GridLayout(tokens)
    renderer = Renderer(tokens, lang=lang)
    single_content = {"slides": [slide_spec]}
    compile_slides(single_content, renderer, grid, tokens)
    # Skip page numbers on the single-slide scratch deck
    fd, tmp = tempfile.mkstemp(prefix="ppt_slide_", suffix=".pptx")
    os.close(fd)
    renderer.save(tmp)
    return tmp


def update_slide(
    pptx_path: str,
    slide_index: int,
    new_slide_spec: dict,
    output_path: Optional[str] = None,
    theme: Union[str, dict] = "academic",
    lang: str = "cn",
) -> str:
    """Replace a slide at the given 1-based index in an existing .pptx.

    替换已有 .pptx 中指定序号（从 1 开始）的幻灯片。先以 _render_single_slide 渲染新 slide
    为临时 deck，再通过 XML 操作替换目标 slide 的 sldId 关系；output_path 为 None 时覆盖原文件。"""
    from pptx import Presentation as _Prs
    from pptx.oxml.ns import qn as _qn
    import copy as _copy
    import shutil

    src = os.path.abspath(pptx_path)
    dst = os.path.abspath(output_path or pptx_path)

    # Render new single slide
    tmp_new = _render_single_slide(new_slide_spec, dst, theme=theme, lang=lang)
    try:
        src_prs = _Prs(src)
        new_prs = _Prs(tmp_new)
        if not 1 <= slide_index <= len(src_prs.slides):
            raise IndexError(
                f"slide_index {slide_index} 超出页码范围 [1, {len(src_prs.slides)}]。修复建议：请传入 1 到 {len(src_prs.slides)} 之间的页码"
            )

        target_slide = src_prs.slides[slide_index - 1]
        new_slide = new_prs.slides[0]

        # Replace shapes on target slide with new slide's shapes via XML
        spTree = target_slide.shapes._spTree
        # Remove existing shapes (but keep nvGrpSpPr/grpSpPr placeholders)
        for child in list(spTree):
            tag = child.tag
            if tag.endswith("}sp") or tag.endswith("}pic") or tag.endswith("}graphicFrame") or tag.endswith("}grpSp") or tag.endswith("}cxnSp"):
                spTree.remove(child)
        # Copy new shapes into spTree
        new_spTree = new_slide.shapes._spTree
        for child in new_spTree:
            tag = child.tag
            if tag.endswith("}sp") or tag.endswith("}pic") or tag.endswith("}graphicFrame") or tag.endswith("}grpSp") or tag.endswith("}cxnSp"):
                spTree.append(_copy.deepcopy(child))

        # Also handle notes (best-effort)
        if new_slide.has_notes_slide:
            new_notes_text = new_slide.notes_slide.notes_text_frame.text
            try:
                if not target_slide.has_notes_slide:
                    target_slide.notes_slide  # trigger creation
                target_slide.notes_slide.notes_text_frame.text = new_notes_text
            except Exception:
                pass

        # Update background if needed
        # (we leave background as-is to preserve theme context)

        src_prs.save(dst)
        return dst
    finally:
        try:
            os.unlink(tmp_new)
        except Exception:
            pass


def add_slide(
    pptx_path: str,
    slide_spec: dict,
    position: Optional[int] = None,
    output_path: Optional[str] = None,
    theme: Union[str, dict] = "academic",
    lang: str = "cn",
) -> str:
    """Insert a new slide at the given position (None = append at end).

    向已有 .pptx 插入一张新幻灯片。position 为 0-based 插入位置，None 表示追加到末尾；
    通过复制 _render_single_slide 产出的临时 deck 的 slide part 到目标 presentation 实现。"""
    from pptx import Presentation as _Prs
    import copy as _copy

    src = os.path.abspath(pptx_path)
    dst = os.path.abspath(output_path or pptx_path)
    prs = _Prs(src)

    # If position is negative or not given, treat as append (or insert before last if -1)
    if position is None:
        idx = len(prs.slides)  # append
    elif position < 0:
        idx = max(0, len(prs.slides) + position + 1)
    else:
        idx = max(0, min(position - 1, len(prs.slides)))

    # Render the new slide to a temp deck so we can steal its shapes and rels
    tmp_new = _render_single_slide(slide_spec, dst, theme=theme, lang=lang)
    try:
        new_prs = _Prs(tmp_new)
        new_slide = new_prs.slides[0]

        # Create new blank slide in src
        blank_layout = prs.slide_layouts[6]
        new_s = prs.slides.add_slide(blank_layout)
        # Background: try to match by filling white; users can regenerate to get theme bg
        from pptx.dml.color import RGBColor
        try:
            bg = new_s.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        except Exception:
            pass

        # Copy shapes XML from new_slide to new_s
        spTree = new_s.shapes._spTree
        new_spTree = new_slide.shapes._spTree
        for child in new_spTree:
            tag = child.tag
            if tag.endswith("}sp") or tag.endswith("}pic") or tag.endswith("}graphicFrame") or tag.endswith("}grpSp") or tag.endswith("}cxnSp"):
                spTree.append(_copy.deepcopy(child))

        # Move newly added slide from last position to position idx
        sldIdLst = prs.slides._sldIdLst
        sldIds = list(sldIdLst)
        new_sldId = sldIds[-1]
        sldIdLst.remove(new_sldId)
        if idx >= len(sldIds) - 1:
            sldIdLst.append(new_sldId)
        else:
            # Insert before the sldId currently at position idx
            ref = sldIds[idx]
            ref.addprevious(new_sldId)

        prs.save(dst)
        return dst
    finally:
        try:
            os.unlink(tmp_new)
        except Exception:
            pass


def delete_slide(
    pptx_path: str,
    slide_index: int,
    output_path: Optional[str] = None,
) -> str:
    """Delete the slide at the given 1-based index from a .pptx.

    删除已有 .pptx 中指定序号（从 1 开始）的幻灯片。通过移除 sldId 与 rId 关系实现；
    python-pptx 会自动维护内部顺序一致性。"""
    from pptx import Presentation as _Prs

    src = os.path.abspath(pptx_path)
    dst = os.path.abspath(output_path or pptx_path)
    prs = _Prs(src)

    if not 1 <= slide_index <= len(prs.slides):
        raise IndexError(
            f"slide_index {slide_index} 超出页码范围 [1, {len(prs.slides)}]。修复建议：请传入 1 到 {len(prs.slides)} 之间的页码"
        )

    # Community-accepted pattern for deleting slides in python-pptx
    slide = prs.slides[slide_index - 1]
    rId = prs.slides._sldIdLst[slide_index - 1].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    # Drop relationship and sldId entry
    prs.part.drop_rel(rId)
    sldIdLst = prs.slides._sldIdLst
    sldIdLst.remove(sldIdLst[slide_index - 1])

    prs.save(dst)
    _relink_slide_order(dst, slide_index)
    return dst


_PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
_DOC_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_SLIDE_REL_TYPE = f"{_DOC_REL_NS}/slide"


def _relink_slide_order(pptx_path: str, deleted_slide_index: int) -> None:
    """Repair internal slide-jump relationships after deleting a slide.

    Any jump target still pointing at the removed slide is redirected to the
    next surviving slide at the same display position, or the previous slide
    when the removed slide was the last one.
    """
    with zipfile.ZipFile(pptx_path, "r") as zin:
        presentation_xml = zin.read("ppt/presentation.xml")
        presentation_rels_xml = zin.read("ppt/_rels/presentation.xml.rels")
        ordered_parts = _ordered_slide_parts(presentation_xml, presentation_rels_xml)
        if not ordered_parts:
            return

        replacement_part = _replacement_slide_part(ordered_parts, deleted_slide_index)
        tmp_path = pptx_path + ".relink_tmp"
        with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.startswith("ppt/slides/_rels/slide") and item.filename.endswith(".xml.rels"):
                    data = _rewrite_slide_jump_relationships(data, ordered_parts, replacement_part)
                zout.writestr(item, data)

    os.replace(tmp_path, pptx_path)


def _ordered_slide_parts(presentation_xml: bytes, presentation_rels_xml: bytes) -> list[str]:
    rels_root = etree.fromstring(presentation_rels_xml)
    rid_to_target = {}
    for rel in rels_root.findall(f"{{{_PKG_REL_NS}}}Relationship"):
        if rel.get("Type") == _SLIDE_REL_TYPE:
            rid_to_target[rel.get("Id")] = _slide_part_name(rel.get("Target", ""))

    pres_root = etree.fromstring(presentation_xml)
    parts: list[str] = []
    for sld_id in pres_root.findall(f".//{{{_PML_NS}}}sldId"):
        rid = sld_id.get(f"{{{_DOC_REL_NS}}}id")
        target = rid_to_target.get(rid)
        if target:
            parts.append(target)
    return parts


def _replacement_slide_part(ordered_parts: list[str], deleted_slide_index: int) -> str:
    replacement_index = min(max(deleted_slide_index - 1, 0), len(ordered_parts) - 1)
    return ordered_parts[replacement_index]


def _slide_part_name(target: str) -> str:
    normalized = (target or "").replace("\\", "/")
    return posixpath.basename(normalized)


def _rewrite_slide_jump_relationships(
    rels_xml: bytes,
    ordered_parts: list[str],
    replacement_part: str,
) -> bytes:
    root = etree.fromstring(rels_xml)
    for rel in root.findall(f"{{{_PKG_REL_NS}}}Relationship"):
        if rel.get("Type") != _SLIDE_REL_TYPE:
            continue
        target = rel.get("Target", "")
        part_name = _slide_part_name(target)
        if part_name in ordered_parts:
            continue
        target_dir = posixpath.dirname((target or "").replace("\\", "/"))
        new_target = replacement_part if target_dir in {"", "."} else f"{target_dir}/{replacement_part}"
        rel.set("Target", new_target)

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


# ----------------------------------------------------------------------
# v1.3 PPT-P0-2: AI image prompt collection helper
# ----------------------------------------------------------------------

def collect_image_prompts(content: dict) -> list[dict]:
    """Walk content and collect all image_prompt fields that have no image_path.

    Returns list of dicts::

        {"page_index": int (1-based), "field": str, "prompt": str,
         "suggested_path": str (relative path hint)}
    """
    if not isinstance(content, dict) or "slides" not in content:
        raise ValidationError("content 必须是包含 'slides' 列表的 dict。修复建议：请传入 {\"slides\": [...]} 结构的对象")
    results = []
    for i, slide in enumerate(content["slides"]):
        def _consider(field: str, obj: dict):
            prompt = obj.get("image_prompt")
            path = obj.get("image_path")
            if prompt and not path:
                results.append({
                    "page_index": i + 1,
                    "field": field,
                    "prompt": prompt,
                    "suggested_path": f"images/slide_{i+1}_{field}.png",
                })
        _consider("image", slide)
        # two_column images (if any)
        for side in ("left", "right"):
            col = slide.get(side)
            if isinstance(col, dict):
                _consider(f"{side}.image", col)
    return results


# ----------------------------------------------------------------------
# v1.3 PPT-P0-3: RUBRIC quality check
# ----------------------------------------------------------------------

def quality_check(content: dict, theme: str = "academic", lang: str = "cn") -> dict:
    """Run RUBRIC quality self-check on PPT semantic content.

    Returns total_score (0-100), passed (>=75), per-dimension scores and issues.
    """
    try:
        from ..shared.quality import (
            contrast_ratio, hex_to_rgb, count_words, count_chars_cjk,
            is_assertive_title, assertive_title_suggestion,
            grade_density, weighted_score, text_length_cjk_aware,
        )
    except ImportError:
        try:
            from shared.quality import (
                contrast_ratio, hex_to_rgb, count_words, count_chars_cjk,
                is_assertive_title, assertive_title_suggestion,
                grade_density, weighted_score, text_length_cjk_aware,
            )
        except ImportError:
            try:
                from skills.shared.quality import (
                    contrast_ratio, hex_to_rgb, count_words, count_chars_cjk,
                    is_assertive_title, assertive_title_suggestion,
                    grade_density, weighted_score, text_length_cjk_aware,
                )
            except ImportError:
                return {
                    "total_score": 0, "passed": False,
                    "error": "shared.quality module not available",
                    "dimensions": {},
                }

    tokens = _resolve_tokens(theme, lang)
    grid = GridLayout(tokens)

    # Resolve theme colors for contrast checks
    def _color_hex(c) -> str:
        # pptx RGBColor is iterable as (r,g,b) ints
        try:
            from pptx.dml.color import RGBColor as _RC
            if isinstance(c, _RC):
                return "#{:02X}{:02X}{:02X}".format(c[0], c[1], c[2])
        except Exception:
            pass
        if isinstance(c, str) and c.startswith("#") and len(c) == 7:
            return c
        if hasattr(c, "__iter__") and not isinstance(c, str):
            try:
                t = tuple(int(v) for v in list(c)[:3])
                return "#{:02X}{:02X}{:02X}".format(*t)
            except Exception:
                return None
        return None

    bg_hex = _color_hex(tokens["color"].get("bg")) or "#FFFFFF"
    title_hex = _color_hex(tokens["color"].get("primary")) or "#1F3864"
    text_hex = _color_hex(tokens["color"].get("text")) or "#333333"
    bg_rgb = hex_to_rgb(bg_hex)
    title_rgb = hex_to_rgb(title_hex)
    text_rgb = hex_to_rgb(text_hex)

    dimensions = {
        "information_density": {"score": 100, "issues": []},
        "title_assertiveness": {"score": 100, "issues": []},
        "visual_balance": {"score": 100, "issues": []},
        "color_contrast": {"score": 100, "issues": []},
        "consistency": {"score": 100, "issues": []},
        "word_count": {"score": 100, "issues": []},
    }

    slides = content.get("slides", [])
    layout_seen = set()
    for i, slide in enumerate(slides):
        page_no = i + 1
        layout = slide.get("layout", "content")
        layout_seen.add(layout)

        # information_density
        d = _density_score(slide, tokens, grid)
        if d == "high":
            dimensions["information_density"]["score"] -= 12
            dimensions["information_density"]["issues"].append({
                "page": page_no,
                "problem": f"信息密度过高（{d}），建议拆分或精简",
                "suggestion": "拆分内容或减少bullets/图表元素",
            })
        elif d == "low" and layout not in ("cover","section","thanks","quote","full_image"):
            dimensions["information_density"]["score"] -= 5
            dimensions["information_density"]["issues"].append({
                "page": page_no,
                "problem": f"信息密度过低（{d}），页面空",
                "suggestion": "合并或补充内容",
            })

        # title_assertiveness
        title_field = {"cover":"title","toc":"title","section":"title",
                       "content":"title","two_column":"title","image_text":"title",
                       "full_image":"title","table":"title","chart":"title",
                       "kpi":"title","quote":None,"summary":"title",
                       "thanks":"title","timeline":"title"}.get(layout, "title")
        title_text = slide.get(title_field) if title_field else None
        if title_text and layout not in ("cover", "toc", "thanks"):
            if not is_assertive_title(title_text):
                dimensions["title_assertiveness"]["score"] -= 8
                dimensions["title_assertiveness"]["issues"].append({
                    "page": page_no,
                    "current_title": title_text,
                    "suggestion": assertive_title_suggestion(title_text),
                })
        # cover subtitle should be assertive/conclusion-ish
        if layout == "cover":
            sub = slide.get("subtitle", "")
            if sub and not is_assertive_title(sub) and len(sub) > 4:
                # soft penalty
                dimensions["title_assertiveness"]["score"] -= 3

        # visual_balance
        if layout == "two_column":
            lb = ((slide.get("left") or {}).get("bullets") or [])
            rb = ((slide.get("right") or {}).get("bullets") or [])
            lc = len(lb) + (1 if (slide.get("left") or {}).get("chart") else 0)
            rc = len(rb) + (1 if (slide.get("right") or {}).get("chart") else 0)
            if abs(lc - rc) >= 4:
                dimensions["visual_balance"]["score"] -= 10
                dimensions["visual_balance"]["issues"].append({
                    "page": page_no,
                    "problem": f"左右栏严重失衡（{lc} vs {rc}）",
                    "suggestion": "调整双栏内容分布",
                })
        if layout == "chart":
            if not slide.get("bullets"):
                dimensions["visual_balance"]["score"] -= 4
                dimensions["visual_balance"]["issues"].append({
                    "page": page_no,
                    "problem": "纯图表页缺少关键结论 bullets",
                    "suggestion": "添加 1–3 条结论要点",
                })
        if layout == "full_image":
            if not slide.get("overlay") and not slide.get("title"):
                dimensions["visual_balance"]["score"] -= 3

        # color_contrast — use theme-level ratios (simplified)
        title_ratio = contrast_ratio(title_rgb, bg_rgb)
        body_ratio = contrast_ratio(text_rgb, bg_rgb)
        if title_ratio < 3.0:
            dimensions["color_contrast"]["score"] -= 10
            dimensions["color_contrast"]["issues"].append({
                "page": page_no, "problem": "标题色对比度不足",
                "ratio": round(title_ratio, 2),
                "suggestion": "加深标题色或提亮背景（WCAG AA 大标题 ≥3:1）",
            })
        if body_ratio < 4.5:
            dimensions["color_contrast"]["score"] -= 10
            dimensions["color_contrast"]["issues"].append({
                "page": page_no, "problem": "正文色对比度不足",
                "ratio": round(body_ratio, 2),
                "suggestion": "加深正文字色（WCAG AA 正文 ≥4.5:1）",
            })

        # word_count
        bullets = slide.get("bullets") or []
        if layout == "two_column":
            for side in ("left","right"):
                col = slide.get(side) or {}
                bullets = bullets + (col.get("bullets") or [])
        if len(bullets) > 6:
            dimensions["word_count"]["score"] -= 8
            dimensions["word_count"]["issues"].append({
                "page": page_no, "bullets": len(bullets),
                "suggestion": f"每页 bullets ≤6 条（当前 {len(bullets)} 条）",
            })
        for bi, b in enumerate(bullets):
            t = b if isinstance(b, str) else (b.get("text") if isinstance(b, dict) else "")
            l = text_length_cjk_aware(t)
            if l > 30:
                dimensions["word_count"]["score"] -= 3
                dimensions["word_count"]["issues"].append({
                    "page": page_no, "bullet_index": bi+1, "length": l,
                    "suggestion": f"单 bullet ≤30 字（当前 {l} 字）：{t[:20]}...",
                })

        # consistency: layout mix looks reasonable (don't over-penalize)
        # (done after loop)

    # consistency: check theme usage across slides (layout jumps etc.)
    # simple: cover should be first, thanks last if present
    if slides:
        if slides[0].get("layout") not in ("cover", "title"):
            dimensions["consistency"]["score"] -= 6
            dimensions["consistency"]["issues"].append({
                "page": 1, "problem": "建议第一页为 cover（封面）",
            })
        if slides[-1].get("layout") not in ("thanks", "summary"):
            dimensions["consistency"]["score"] -= 3
            dimensions["consistency"]["issues"].append({
                "page": len(slides), "problem": "建议最后一页为 thanks（致谢）或 summary（总结）",
            })

    # clamp scores to 0..100
    for k in dimensions:
        dimensions[k]["score"] = max(0, min(100, int(dimensions[k]["score"])))

    weights = {
        "information_density": 25,
        "title_assertiveness": 20,
        "visual_balance": 15,
        "color_contrast": 15,
        "consistency": 15,
        "word_count": 10,
    }
    total, _ = weighted_score(dimensions, weights)
    return {
        "total_score": total,
        "passed": total >= 75,
        "dimensions": dimensions,
        "weights": weights,
    }


# ----------------------------------------------------------------------
# v1.5: Taste / craft preflight
# ----------------------------------------------------------------------
def taste_check(
    content: dict,
    theme: str = "academic",
    lang: str = "cn",
    strict: bool = False,
) -> dict:
    """Run an Impeccable/Taste-inspired deck craft preflight.

    This complements ``quality_check`` by detecting templated AI copy,
    placeholder text, weak visual intent, and layout repetition before the deck
    is rendered.
    """
    from .taste import taste_check as _taste_check

    return _taste_check(content, theme=theme, lang=lang, strict=strict)


def visual_redesign_check(content: dict, strict: bool = False) -> dict:
    """Run anti-flat infographic-structure checks before rendering."""
    from .visual_redesign import visual_redesign_check as _check

    return _check(content, strict=strict)


def build_visual_redesign_prompt(content: dict, page_count: int | None = None) -> str:
    """Build a visual-redesign planning prompt for PPT generation."""
    from .visual_redesign import build_visual_redesign_prompt as _build

    return _build(content, page_count=page_count)


def apply_visual_redesign_guidance(content: dict) -> dict:
    """Attach visual_plan guidance to core narrative slides."""
    from .visual_redesign import apply_visual_redesign_guidance as _apply

    return _apply(content)


# ----------------------------------------------------------------------
# v1.3 PPT-P0-4: Narrative / story check
# ----------------------------------------------------------------------

def assertive_title(topic: str, data_point: str = "") -> str:
    """Generate an assertive title suggestion from topic + optional data point.

    Pure heuristic (no LLM).
    """
    topic = (topic or "").strip()
    if not topic:
        return "请提供主题以生成断言式标题"
    if data_point:
        templates = [
            "{topic}达{data_point}，显著超预期",
            "{topic}突破{data_point}，行业领先",
            "{topic}实现{data_point}的核心增长",
            "{topic}提升至{data_point}，验证方案有效",
        ]
    else:
        templates = [
            "{topic}实现显著提升",
            "{topic}驱动核心增长",
            "{topic}验证方案有效",
            "{topic}达成关键目标",
        ]
    import random
    return random.choice(templates).format(topic=topic, data_point=data_point)


_SCQA_KEYWORDS = {
    "S": ("背景","现状","情境","概况","s","situation","context"),
    "C": ("挑战","问题","冲突","痛点","矛盾","c","complication","challenge"),
    "Q": ("分析","关键问题","问题分析","核心问题","q","question"),
    "A": ("对策","建议","方案","结论","对策建议","a","answer","resolution"),
}


def story_check(content: dict, method: str = "pyramid") -> dict:
    """Check narrative structure of the deck.

    Methods: pyramid (SCQA-like conclusion-first), scqa, chronological.
    """
    if not isinstance(content, dict) or "slides" not in content:
        raise ValidationError("content 必须是包含 'slides' 列表的 dict。修复建议：请传入 {\"slides\": [...]} 结构的对象")
    try:
        from ..shared.quality import is_assertive_title
    except ImportError:
        try:
            from shared.quality import is_assertive_title
        except ImportError:
            try:
                from skills.shared.quality import is_assertive_title
            except ImportError:
                is_assertive_title = lambda t: bool(t)

    slides = content["slides"]
    issues: list[dict] = []
    compliant = True

    if method == "pyramid":
        # Cover subtitle should be conclusion
        if slides and slides[0].get("layout") == "cover":
            sub = slides[0].get("subtitle") or ""
            if sub and not is_assertive_title(sub):
                compliant = False
                issues.append({
                    "page_index": 1, "field": "subtitle",
                    "problem": "封面副标题应是核心结论句",
                    "suggestion": assertive_title(sub or "核心观点"),
                })
        # Chapters 3-5 via section pages
        sections = [s for s in slides if s.get("layout") == "section"]
        if not (3 <= len(sections) <= 5):
            issues.append({
                "page_index": 0, "field": "sections",
                "problem": f"章节页数量 {len(sections)}，建议3–5个以符合MECE",
                "suggestion": "拆分为3–5个章节页，做到相互独立、完全穷尽",
            })
            if len(sections) < 2:
                compliant = False
        # Every non-cover/section/thanks title should be assertive
        for i, s in enumerate(slides):
            layout = s.get("layout")
            if layout in ("cover","section","thanks","toc","quote"):
                continue
            t = s.get("title") or ""
            if t and not is_assertive_title(t):
                issues.append({
                    "page_index": i+1, "field": "title",
                    "problem": f"标题「{t}」为名词短语，建议改为断言式",
                    "suggestion": assertive_title(t),
                })
        # Summary should echo cover conclusion
        summary_pages = [s for s in slides if s.get("layout") == "summary"]
        if not summary_pages:
            issues.append({
                "page_index": len(slides), "field": "summary",
                "problem": "缺少 summary（总结回顾）页",
                "suggestion": "添加 summary 页回扣封面核心结论",
            })

    elif method == "scqa":
        # Check S/C/Q/A chapter presence by title keywords
        buckets = {k: [] for k in "SCQA"}
        for i, s in enumerate(slides):
            title = s.get("title") or ""
            if s.get("layout") == "section":
                for k, kws in _SCQA_KEYWORDS.items():
                    if any(kw in title for kw in kws):
                        buckets[k].append(i+1); break
        missing = [k for k, v in buckets.items() if not v]
        if missing:
            compliant = False
            mapping = {"S":"情境(Situation)","C":"冲突(Complication)","Q":"问题(Question)","A":"答案(Answer)"}
            issues.append({
                "page_index": 0, "field": "sections",
                "problem": f"SCQA 缺失章节: {', '.join(mapping[m] for m in missing)}",
                "suggestion": "按 背景→挑战→关键问题→对策建议 组织章节页",
            })
        # Final summary = A
        if slides and slides[-1].get("layout") not in ("summary","thanks"):
            issues.append({
                "page_index": len(slides), "field": "last",
                "problem": "最后一页应是 summary/Answer 或 thanks",
            })

    elif method == "chronological":
        # timeline events date-sorted
        import re
        date_re = re.compile(r"(\d{4})[-/.年]?(\d{1,2})?")
        for i, s in enumerate(slides):
            if s.get("layout") != "timeline":
                continue
            events = s.get("events") or []
            dates = []
            for j, ev in enumerate(events):
                if not isinstance(ev, dict):
                    continue
                d = ev.get("date","")
                m = date_re.search(str(d))
                if m:
                    y = int(m.group(1)); mo = int(m.group(2) or 1)
                    dates.append((y*12+mo, j, d))
                else:
                    issues.append({
                        "page_index": i+1, "field": f"events[{j}].date",
                        "problem": f"事件日期无法解析: {d}",
                        "suggestion": "使用YYYY-MM或YYYY年M月格式",
                    })
            if len(dates) >= 2:
                ordered = sorted(dates)
                if [d[0] for d in dates] != [d[0] for d in ordered]:
                    compliant = False
                    issues.append({
                        "page_index": i+1, "field": "events",
                        "problem": "timeline 事件未按时间升序排列",
                        "suggestion": "将事件按时间顺序重新排列",
                    })
        if not any(s.get("layout") == "timeline" for s in slides):
            issues.append({
                "page_index": 0, "field": "timeline",
                "problem": "时间顺序叙事应包含至少一页 timeline",
            })
    else:
        return {
            "method": method, "compliant": False,
            "issues": [{"page_index":0,"field":"method",
                        "problem": f"未知叙事方法 {method}",
                        "suggestion": "选择 pyramid / scqa / chronological"}],
            "suggested_outline": None,
        }

    if issues and compliant:
        # if there are only soft suggestions, keep compliant=True if no critical
        critical = sum(1 for i in issues if i.get("page_index",0)>0)
        if critical > 4:
            compliant = False

    # simple suggested outline (none unless severely broken)
    suggested_outline = None
    if method == "pyramid" and not compliant:
        cover = slides[0] if slides else {}
        suggested_outline = [
            {"layout":"cover","title":cover.get("title","主题"),
             "subtitle": assertive_title(cover.get("title","主题"))},
            {"layout":"toc","title":"目录"},
            {"layout":"section","number":"01","title":"背景与现状"},
            {"layout":"content","title":assertive_title("现状与数据"),"bullets":["要点1","要点2"]},
            {"layout":"section","number":"02","title":"核心发现"},
            {"layout":"content","title":assertive_title("发现"),"bullets":["发现1","发现2"]},
            {"layout":"section","number":"03","title":"方案与建议"},
            {"layout":"content","title":assertive_title("建议"),"bullets":["建议1","建议2"]},
            {"layout":"summary","title":"总结"},
            {"layout":"thanks","title":"感谢聆听"},
        ]

    return {
        "method": method,
        "compliant": compliant,
        "issues": issues,
        "suggested_outline": suggested_outline,
    }


# ----------------------------------------------------------------------
# v1.4 P2-1: Mini (sparkline) chart SVG helpers
# ----------------------------------------------------------------------
def _load_shared_module(module_name: str):
    return import_shared(module_name)


def _mini_theme_colors(theme: str) -> dict:
    shared_themes = _load_shared_module("themes")
    tc = dict(shared_themes.get_theme_palette(theme))
    tc["_chart_palette"] = shared_themes.get_chart_palette(theme)
    return tc


def mini_bar(data, width: int = 120, height: int = 40, color: Optional[str] = None,
             *, theme: str = "academic", **kwargs) -> str:
    """Render an SVG mini bar-chart string for inline placement.

    生成 SVG 迷你柱状图字符串，可传入 KPI 页的 mini_charts 字段或任何接受 SVG 的
    渲染调用（如 free_shape）。data 接受数值列表；width/height 单位像素；color 为
    None 时自动从主题 chart 调色板取色。"""
    _mb = _load_shared_module("svg_engine").mini_bar
    tc = _mini_theme_colors(theme)
    return _mb(data, width=width, height=height, color=color, theme_colors=tc, **kwargs)


def mini_line(data, width: int = 120, height: int = 40, color: Optional[str] = None,
              *, theme: str = "academic", fill: bool = True, **kwargs) -> str:
    """Render an SVG mini line-chart string for inline placement.

    生成 SVG 迷你折线图字符串，参数同 mini_bar；fill=True 时绘制折线下方半透明填充。
    返回的 SVG 字符串可直接嵌入 KPI 指标卡或任意 slide 节点的 svg 字段。"""
    _ml = _load_shared_module("svg_engine").mini_line
    tc = _mini_theme_colors(theme)
    return _ml(data, width=width, height=height, color=color, theme_colors=tc, fill=fill, **kwargs)


def mini_pie(parts, size: int = 80, colors=None, *, theme: str = "academic", **kwargs) -> str:
    """Render an SVG mini pie-chart string for inline placement.

    生成 SVG 迷你饼图字符串。parts 为数值列表（按比例划分扇区），colors 可指定扇区颜色
    列表；默认从主题 4 色 chart palette 循环取色。适合 KPI 页占比展示。"""
    _mp = _load_shared_module("svg_engine").mini_pie
    tc = _mini_theme_colors(theme)
    return _mp(parts, size=size, colors=colors, theme_colors=tc, **kwargs)


# ----------------------------------------------------------------------
# v1.6.3: A-6 模板匹配评分 (PPT 端)
# ----------------------------------------------------------------------
def check_template_match(
    tokens: dict,
    template_name: Optional[str] = None,
    extracted: Optional[dict] = None,
) -> dict:
    """PPT 端轻量模板匹配评分（0~100）。

    评分维度：
      - 颜色对比度：text vs bg（WCAG AA ≥4.5）、title_bar_bg vs title 文字（≥4.5）
      - 字体完整性：family.heading/body/default 都有
      - 模板一致性：若传 template_name + extracted，验证 tokens.primary 与模板主题色协调

    Returns:
        { "match_score": int 0-100, "passed": bool (≥60), "warnings": list[str], "details": dict }
    """
    warnings: list[str] = []
    score = 100
    body_ratio = None
    title_ratio = None

    color = tokens.get("color", {}) or {}
    bg = color.get("bg", "#FFFFFF")
    text = color.get("text", "#333333")
    primary = color.get("primary", "#1F3864")
    title_bar_bg = color.get("title_bar_bg", primary)
    title_color = color.get("title", primary)

    (hex_to_rgb,) = import_shared("color_palette", attrs=["hex_to_rgb"])
    (contrast_ratio,) = import_shared("quality", attrs=["contrast_ratio"])

    # 1) 正文对比度
    try:
        body_ratio = contrast_ratio(hex_to_rgb(text), hex_to_rgb(bg))
        if body_ratio < 4.5:
            score -= 20
            warnings.append(f"正文对比度 {body_ratio:.2f}:1 低于 WCAG AA 4.5:1")
    except Exception:
        score -= 5
        warnings.append("正文对比度计算失败")

    # 2) 标题栏对比度
    try:
        title_ratio = contrast_ratio(hex_to_rgb(title_color), hex_to_rgb(title_bar_bg))
        if title_ratio < 4.5:
            score -= 15
            warnings.append(f"标题栏对比度 {title_ratio:.2f}:1 低于 WCAG AA 4.5:1")
    except Exception:
        pass

    # 3) 字体完整性
    family = (tokens.get("font", {}) or {}).get("family", {}) or {}
    for key in ("heading", "body", "default"):
        if not family.get(key):
            score -= 5
            warnings.append(f"字体缺少 {key}")

    # 4) 模板一致性
    if template_name and extracted:
        ext_tokens = getattr(extracted, "tokens", None) or extracted.get("tokens", {}) if isinstance(extracted, dict) else {}
        ext_color = (ext_tokens or {}).get("color", {}) or {}
        ext_primary = (ext_color.get("primary") or "").upper()
        cur_primary = str(primary).upper()
        if ext_primary and cur_primary and ext_primary != cur_primary:
            warnings.append(
                f"当前 primary {cur_primary} 与模板 {template_name} 的 {ext_primary} 不一致"
            )
            score -= 3

    score = max(0, min(100, score))
    return {
        "match_score": score,
        "passed": score >= 60,
        "warnings": warnings,
        "details": {
            "body_contrast": body_ratio,
            "title_contrast": title_ratio,
        },
    }


def mini_chart_svg(kind: str, data, *, theme: str = "academic", **kwargs) -> str:
    """Dispatch helper returning SVG for bar/line/pie mini charts by kind.

    统一入口的迷你图表 SVG 生成函数。kind 取 ``"bar"``/``"line"``/``"pie"``，
    其余 kwargs 透传给对应 mini_* 函数；自动注入主题色板，调用方无需手动传 theme_colors。"""
    _mc = _load_shared_module("svg_engine").mini_chart
    tc = _mini_theme_colors(theme)
    kwargs.setdefault("theme_colors", tc)
    return _mc(kind, data, **kwargs)
